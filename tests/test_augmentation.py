"""Preservation, input integrity, and whole-run determinism for P65-P71."""
from __future__ import annotations

import json
import os
import subprocess
import sys
from pathlib import Path

import pytest

from proofs.augment import Inputs, ROOT, canonical, digest, layout_graph, render_page, check_rendered


@pytest.fixture(scope="module")
def runs(tmp_path_factory):
    base = tmp_path_factory.mktemp("augmentation")
    outputs = []
    for seed in ("1", "424242"):
        out = base / seed
        completed = subprocess.run([sys.executable, "-m", "proofs.augment", "--output", str(out)],
                                   cwd=ROOT, env={**os.environ, "PYTHONHASHSEED": seed}, capture_output=True, text=True, timeout=90)
        assert completed.returncode == 0, completed.stderr
        outputs.append(out)
    return outputs


def test_whole_pipeline_is_identical_across_processes_and_hash_seeds(runs):
    a, b = runs
    assert sorted(p.name for p in a.iterdir()) == sorted(p.name for p in b.iterdir())
    for p in a.iterdir():
        assert p.read_bytes() == (b / p.name).read_bytes(), p.name
    report = json.loads((a / "verification.json").read_text())
    assert report["runtime"] == {"external_execution_attempts": [], "model_import_attempts": []}


def test_all_policy_inputs_and_meetings_are_preserved(runs):
    model = json.loads((runs[0] / "model.json").read_text())
    policy = model["policy"]
    assert policy["checks"] == {"zen_table_agreement": 19, "dmn_table_agreement": 19, "independent_readers_agree": True}
    assert len(policy["meetings"]) == 8
    assert policy["threshold"] == 10 and policy["total"] == 18


def test_graph_ids_labels_and_directed_edges_are_preserved(runs):
    model = json.loads((runs[0] / "model.json").read_text())
    for graph in model["processes"]:
        source = json.loads((ROOT / graph["source"]).read_text())
        assert [n["id"] for n in graph["nodes"]] == source["order"]
        assert graph["edges"] == source["forced"]
        assert len({(n["x"], n["y"]) for n in graph["nodes"]}) == len(graph["nodes"])
        by_id = {n["id"]: n for n in graph["nodes"]}
        assert all(by_id[a]["x"] < by_id[b]["x"] for a, b in graph["edges"])


def test_office_values_text_and_connectors_survive(runs):
    from openpyxl import load_workbook
    from docx import Document
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    model = json.loads((runs[0] / "model.json").read_text())
    for table in model["records"]:
        workbook = load_workbook(ROOT / table["source"], data_only=False)
        assert [table["columns"], *table["rows"]] == [list(r) for r in workbook.active.values]
    for document in model["documents"]:
        original = Document(ROOT / document["source"])
        assert [b["text"] for b in document["blocks"] if b["kind"] == "paragraph"] == [p.text for p in original.paragraphs]
        assert [b["rows"] for b in document["blocks"] if b["kind"] == "table"] == [[[c.text for c in row.cells] for row in t.rows] for t in original.tables]
    deck = Presentation(ROOT / model["deck"]["source"])
    for view, slide in zip(model["deck"]["slides"], deck.slides, strict=True):
        assert view["headings"] + [n["text"] for n in view["nodes"]] == [s.text for s in slide.shapes if s.has_text_frame]
        expected = [[s._element.xpath(".//a:stCxn")[0].get("id"), s._element.xpath(".//a:endCxn")[0].get("id")] for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE]
        assert view["edges"] == expected


def test_changed_or_unregistered_source_is_rejected(tmp_path):
    (tmp_path / "proofs").mkdir()
    (tmp_path / "proofs/register.json").write_text(json.dumps({"source.txt": digest(b"original")}))
    (tmp_path / "source.txt").write_bytes(b"changed")
    with pytest.raises(ValueError, match="changed"):
        Inputs(tmp_path).read("source.txt")
    with pytest.raises(ValueError, match="unregistered"):
        Inputs(tmp_path).read("proofs/register.json")


def test_runtime_rejects_model_import_and_network():
    for operation in ["import openai", "import socket; socket.create_connection(('example.com',443))"]:
        code = "from proofs.augment import install_boundary; install_boundary(); " + operation
        result = subprocess.run([sys.executable, "-c", code], cwd=ROOT, capture_output=True, text=True, timeout=10)
        assert result.returncode != 0
        assert "outside" in result.stderr and "contract" in result.stderr


def test_untrusted_source_text_cannot_escape_data_or_dom(tmp_path):
    from lxml import html
    config = json.loads((ROOT / "proofs/presentation.json").read_text())
    payload = {"version": 1, "sources": {}, "documents": [{"id":"x", "title":"x", "source":"x", "blocks":[{"kind":"paragraph", "style":"Normal", "text":"</script><img src=x onerror=alert(1)>"}]}]}
    path = tmp_path / "hostile.html"
    path.write_text(render_page(payload, "briefing", config))
    assert check_rendered(path, payload)["model_roundtrip"]
    assert not html.fromstring(path.read_bytes()).xpath("//img")


def test_graph_contract_rejects_cycles_and_unknown_endpoints():
    nodes = [{"id":"a"}, {"id":"b"}]
    with pytest.raises(ValueError, match="acyclic"):
        layout_graph(nodes, [["a","b"],["b","a"]])
    with pytest.raises(ValueError, match="endpoints"):
        layout_graph(nodes, [["a","c"]])
