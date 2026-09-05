"""Presentation chains P65-P71 over registered outputs. No model intervention.

    .venv/bin/python -m proofs.augment --record

Design choices are versioned data. The fixed dispatcher executes every handoff.
Sources are checked against the existing register before any presentation is built.
The originals are read only; new artifacts receive new proof numbers and digests.
"""
from __future__ import annotations

import argparse
import hashlib
import importlib.abc
import importlib.metadata
import io
import json
import re
import sys
import types
from pathlib import Path
from typing import Any

ROOT = Path(__file__).resolve().parents[1]
HERE = Path(__file__).resolve().parent
DEFAULT_OUT = HERE / "out" / "augmentation"
MODEL_CLIENTS = frozenset({"openai", "anthropic", "litellm", "langchain", "langchain_openai", "transformers", "ollama"})


def install_boundary() -> dict[str, Any]:
    """Fail closed on model-client imports, network calls, and child execution."""
    state: dict[str, Any] = {"model_import_attempts": [], "external_execution_attempts": []}

    class BlockModels(importlib.abc.MetaPathFinder):
        def find_spec(self, fullname: str, path: Any = None, target: Any = None) -> Any:
            if fullname.split(".")[0] in MODEL_CLIENTS:
                state["model_import_attempts"].append(fullname)
                raise RuntimeError(f"Model client is outside the execution contract: {fullname}")
            return None

    def audit(event: str, args: Any) -> None:
        if event in {"socket.connect", "socket.getaddrinfo", "subprocess.Popen", "os.system", "os.exec", "os.posix_spawn"}:
            state["external_execution_attempts"].append(event)
            raise RuntimeError(f"External execution is outside the presentation contract: {event}")

    if MODEL_CLIENTS.intersection(sys.modules):
        raise RuntimeError("A model client was loaded before the presentation boundary")
    sys.meta_path.insert(0, BlockModels())
    sys.addaudithook(audit)
    return state


def digest(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def canonical(value: Any) -> bytes:
    return (json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":")) + "\n").encode("utf-8")


class Inputs:
    def __init__(self, root: Path):
        self.root = root
        self.register = json.loads((root / "proofs/register.json").read_text())
        self.used: dict[str, str] = {}

    def read(self, relative: str) -> bytes:
        path = (self.root / relative).resolve()
        if not path.is_relative_to(self.root.resolve()):
            raise ValueError("Source path leaves the repository")
        data = path.read_bytes()
        actual = digest(data)
        expected = self.register.get(relative)
        if expected is None or actual != expected:
            raise ValueError(f"Source is unregistered or changed: {relative}")
        self.used[relative] = actual
        return data


def read_sheet(data: bytes) -> dict[str, Any]:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), data_only=False, read_only=True)
    if len(wb.worksheets) != 1:
        raise ValueError("This view requires exactly one worksheet")
    values = [list(row) for row in wb.active.values]
    wb.close()
    if not values or any(not isinstance(c, str) for c in values[0]):
        raise ValueError("A nonempty string header is required")
    if any(isinstance(v, str) and v.startswith("=") for row in values[1:] for v in row):
        raise ValueError("Formula cells require an explicit evaluated-value adapter")
    return {"columns": values[0], "rows": values[1:]}


def policy_model(inputs: Inputs, spec: dict[str, Any]) -> dict[str, Any]:
    import zen
    from python_calamine import CalamineWorkbook
    from SpiffWorkflow.bpmn.script_engine import PythonScriptEngine
    from SpiffWorkflow.dmn.engine.DMNEngine import DMNEngine
    from SpiffWorkflow.dmn.parser.BpmnDmnParser import BpmnDmnParser

    policy_bytes = inputs.read(spec["policy"])
    policy = json.loads(policy_bytes)
    decision = zen.ZenEngine().create_decision(policy_bytes.decode())
    truth_bytes = inputs.read(spec["truth"])
    truth = read_sheet(truth_bytes)
    if truth["columns"] != ["present_voting_members", "majority_reachable", "rule"]:
        raise ValueError("Unexpected truth table schema")
    independent = CalamineWorkbook.from_filelike(io.BytesIO(truth_bytes)).get_sheet_by_index(0).to_python()
    if independent != [truth["columns"], *truth["rows"]]:
        raise ValueError("Independent workbook readers disagree")
    values = truth["rows"]
    if [r[0] for r in values] != list(range(len(values))):
        raise ValueError("Input space must be the complete consecutive integer interval 0..N")
    parser = BpmnDmnParser()
    parser.add_dmn_str(inputs.read(spec["dmn"]).decode())
    parsed = parser.dmn_parsers["majority"]
    parsed.parse()
    tables = parsed.decision.decision_tables if hasattr(parsed.decision, "decision_tables") else parsed.decision.decisionTables
    dmn = DMNEngine(tables[0])
    for n, reachable, rule in values:
        if type(n) is not int or type(reachable) is not bool or not isinstance(rule, str):
            raise ValueError("Truth table values have invalid types")
        result = decision.evaluate({"present": n})["result"]
        task = types.SimpleNamespace(data={"present": n}, workflow=types.SimpleNamespace(script_engine=PythonScriptEngine()))
        if result != {"majority_reachable": reachable, "rule": rule} or dmn.result(task)["majority_reachable"] != reachable:
            raise ValueError(f"Policy, DMN, and stored table disagree at {n}")
    meetings = read_sheet(inputs.read(spec["meetings"]))
    if meetings["columns"] != ["meeting", "present_voting", "majority_reachable"]:
        raise ValueError("Unexpected meeting schema")
    for _, n, result in meetings["rows"]:
        if type(n) is not int or not 0 <= n < len(values) or values[n][1] != result:
            raise ValueError("Meeting does not match the stored policy table")
    table = next(n["content"] for n in policy["nodes"] if n["type"] == "decisionTableNode")
    rules = {r["_id"]: {"description": r["_description"], "condition": r["i1"]} for r in table["rules"]}
    if any(r[2] not in rules for r in values):
        raise ValueError("Unknown rule reference")
    return {"title": spec["title"], "total": len(values) - 1, "truth": values,
            "threshold": next(r[0] for r in values if r[1]), "meetings": meetings["rows"],
            "rules": rules, "note": spec["roster_note"], "source": spec["policy"],
            "checks": {"zen_table_agreement": len(values), "dmn_table_agreement": len(values), "independent_readers_agree": True}}


def layout_graph(nodes: list[dict[str, Any]], edges: list[list[str]]) -> dict[str, Any]:
    import networkx as nx

    ids = [n["id"] for n in nodes]
    if len(ids) != len(set(ids)) or any(a not in ids or b not in ids for a, b in edges):
        raise ValueError("Graph IDs or endpoints are invalid")
    graph = nx.DiGraph()
    graph.add_nodes_from(ids)
    graph.add_edges_from(edges)
    if not nx.is_directed_acyclic_graph(graph):
        raise ValueError("The dependency view requires an acyclic graph")
    ordinal = {node: i for i, node in enumerate(ids)}
    layers = [sorted(layer, key=ordinal.__getitem__) for layer in nx.topological_generations(graph)]
    positions = {node: {"x": 38 + c * 270, "y": 38 + r * 106} for c, layer in enumerate(layers) for r, node in enumerate(layer)}
    return {"nodes": [{**n, **positions[n["id"]]} for n in nodes], "edges": edges,
            "width": max(540, len(layers) * 270 + 38), "height": max(260, max(map(len, layers), default=0) * 106 + 38),
            "layers": len(layers)}


def process_model(inputs: Inputs, spec: dict[str, Any]) -> dict[str, Any]:
    key = spec["id"]
    order_path = f"proofs/out/P2/{key}/ordered_steps.json"
    facts_path = f"proofs/out/P1/{key}/facts.jsonl"
    ordering = json.loads(inputs.read(order_path))
    if ordering["cycle"]:
        raise ValueError("Recorded ordering contains a cycle")
    facts = [json.loads(row) for row in inputs.read(facts_path).decode().splitlines()]
    events = {f["args"][0]: f for f in facts if f["predicate"] == "event"}
    nodes = [{"id": i, "label": events[i]["args"][1], "quote": events[i]["quote"], "sentence": events[i]["sentence_id"]} for i in ordering["order"]]
    model = layout_graph(nodes, ordering["forced"])
    return {**model, "id": key, "title": spec["title"], "source": order_path}


def records_model(inputs: Inputs, spec: dict[str, Any]) -> dict[str, Any]:
    table = read_sheet(inputs.read(spec["path"]))
    if len(spec["labels"]) != len(table["columns"]):
        raise ValueError("Column labels must map one-to-one to source columns")
    return {**table, **{k: spec[k] for k in ["id", "title", "labels", "description"]}, "source": spec["path"]}


def document_model(inputs: Inputs, spec: dict[str, Any]) -> dict[str, Any]:
    from docx import Document
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(io.BytesIO(inputs.read(spec["path"])))
    blocks = []
    for child in doc.element.body:
        if child.tag == qn("w:p"):
            p = Paragraph(child, doc)
            blocks.append({"kind": "paragraph", "style": p.style.name, "text": p.text})
        elif child.tag == qn("w:tbl"):
            t = Table(child, doc)
            blocks.append({"kind": "table", "rows": [[c.text for c in row.cells] for row in t.rows]})
    return {"id": spec["id"], "title": spec["title"], "source": spec["path"], "blocks": blocks}


def slide_model(inputs: Inputs, spec: dict[str, Any]) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    deck = Presentation(io.BytesIO(inputs.read(spec["path"])))
    slides = []
    for slide in deck.slides:
        nodes, edges, headings = [], [], []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                starts = shape._element.xpath(".//a:stCxn")
                ends = shape._element.xpath(".//a:endCxn")
                if len(starts) != 1 or len(ends) != 1:
                    raise ValueError("A connector must have explicit shape endpoints")
                edges.append([starts[0].get("id"), ends[0].get("id")])
            elif shape.has_text_frame:
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    headings.append(shape.text)
                else:
                    nodes.append({"id": str(shape.shape_id), "label": shape.text.split("\n")[0], "text": shape.text})
            else:
                raise ValueError("This slide reader requires text nodes and explicit connectors")
        slides.append({**layout_graph(nodes, edges), "headings": headings})
    return {"title": spec["title"], "source": spec["path"], "slides": slides}


def render_page(payload: dict[str, Any], initial: str, config: dict[str, Any]) -> str:
    import jinja2

    env = jinja2.Environment(loader=jinja2.FileSystemLoader(str(HERE / "ui")), autoescape=True, undefined=jinja2.StrictUndefined)
    return env.get_template("workbench.html").render(
        model=payload, initial=initial, theme=config["theme"],
        css=(HERE / "ui/workbench.css").read_text(), js=(HERE / "ui/workbench.js").read_text())


def check_rendered(path: Path, payload: dict[str, Any]) -> dict[str, Any]:
    from lxml import html

    doc = html.fromstring(path.read_bytes())
    embedded = json.loads(doc.get_element_by_id("model").text)
    if embedded != payload:
        raise ValueError("Presentation model changed during serialization")
    if doc.get("lang") != "en-US":
        raise ValueError("US English locale is required")
    if doc.xpath("//script[@src] | //link[@href] | //iframe | //img[@src]"):
        raise ValueError("Standalone pages must have no external assets or frames")
    csp = doc.xpath("//meta[@http-equiv='Content-Security-Policy']/@content")
    if not csp or "connect-src 'none'" not in csp[0]:
        raise ValueError("Browser network access must be disabled")
    return {"model_roundtrip": True, "locale": "en-US", "external_assets": 0, "browser_connections": "disabled"}


def build(output: Path, boundary: dict[str, Any]) -> dict[str, Any]:
    config = json.loads((HERE / "presentation.json").read_text())
    if config["runtime"] != {"llm_reasoning": "forbidden", "network": "forbidden"} or config["locale"] != "en-US":
        raise ValueError("Execution contract cannot be relaxed by configuration")
    inputs = Inputs(ROOT)
    artifacts: list[dict[str, Any]] = []
    output.mkdir(parents=True, exist_ok=True)
    payload: dict[str, Any] = {"version": 1, "policy": policy_model(inputs, config["policy"])}
    payload["processes"] = [process_model(inputs, s) for s in config["processes"]]
    payload["records"] = [records_model(inputs, s) for s in config["records"]]
    payload["documents"] = [document_model(inputs, s) for s in config["documents"]]
    payload["deck"] = slide_model(inputs, config["deck"])
    payload["sources"] = inputs.used

    definitions = [
        ("P65", "policy", "policy, truth table, meetings -> checked policy interface", [config["policy"][k] for k in ["policy", "truth", "meetings", "dmn"]],
         "zen.ZenDecision.evaluate; CalamineWorkbook.from_filelike; DMNEngine.result; Jinja2.Template.render",
         [f"all {len(payload['policy']['truth'])} inputs agree across Zen, DMN, and the stored truth table", f"meeting rows {len(payload['policy']['meetings'])}; no formula or policy rewrite"]),
        ("P66", "process", "ordered steps, facts -> layered graph -> process explorer", [f"proofs/out/P2/{s['id']}/ordered_steps.json" for s in config["processes"]] + [f"proofs/out/P1/{s['id']}/facts.jsonl" for s in config["processes"]],
         "networkx.is_directed_acyclic_graph; networkx.topological_generations; Jinja2.Template.render",
         [f"{p['id']}: nodes {len(p['nodes'])}, edges {len(p['edges'])}; IDs, labels, quotes, and edge endpoints preserved" for p in payload["processes"]]),
        ("P67", "records", "workbooks -> typed tables -> searchable records interface", [s["path"] for s in config["records"]],
         "openpyxl.load_workbook; Worksheet.values; Jinja2.Template.render; Array.filter; Array.sort; Array.slice",
         [f"{r['id']}: {len(r['rows'])} rows, {len(r['columns'])} columns; original values and types retained" for r in payload["records"]]),
        ("P68", "briefing", "documents -> ordered blocks -> readable briefings", [s["path"] for s in config["documents"]],
         "docx.Document; Paragraph.text; Table.rows; Jinja2.Template.render",
         [f"{d['id']}: {len(d['blocks'])} blocks; paragraph and table order preserved" for d in payload["documents"]]),
        ("P69", "slides", "deck -> text nodes, connected edges -> readable slide view", [config["deck"]["path"]],
         "pptx.Presentation; slide.shapes; lxml XPath stCxn/endCxn; networkx.topological_generations; Jinja2.Template.render",
         [f"slides {len(payload['deck']['slides'])}; text nodes {sum(len(s['nodes']) for s in payload['deck']['slides'])}; connectors {sum(len(s['edges']) for s in payload['deck']['slides'])}; shape endpoints retained"]),
    ]
    for pid, view, title, sources, functions, shows in definitions:
        path = output / f"{view}.html"
        key = {"policy": "policy", "process": "processes", "records": "records", "briefing": "documents", "slides": "deck"}[view]
        view_payload = {"version": 1, key: payload[key], "sources": {s: inputs.used[s] for s in sources}}
        path.write_text(render_page(view_payload, view, config), encoding="utf-8")
        checks = check_rendered(path, view_payload)
        artifacts.append({"pid": pid, "label": view, "title": title, "sources": sources,
                          "functions": functions, "shows": shows, "checks": checks,
                          "file": path.name, "sha256": digest(path.read_bytes())})
    path = output / "workbench.html"
    path.write_text(render_page(payload, "policy", config), encoding="utf-8")
    artifacts.append({"pid": "P70", "label": "workbench", "title": "checked views -> shared navigation -> standalone workbench",
                      "sources": [a["file"] for a in artifacts], "functions": "Jinja2.Template.render; DOM.addEventListener; DOM.hidden",
                      "shows": ["five views composed from the same model; no external assets", "all interactions use fixed handlers and typed records"],
                      "checks": check_rendered(path, payload), "file": path.name, "sha256": digest(path.read_bytes())})
    (output / "model.json").write_bytes(canonical(payload))
    versions = {n: importlib.metadata.version(n) for n in ["Jinja2", "networkx", "zen-engine", "SpiffWorkflow", "openpyxl", "python-calamine", "python-docx", "python-pptx", "lxml", "PyYAML"]}
    if boundary["model_import_attempts"] or boundary["external_execution_attempts"]:
        raise ValueError("Execution attempted to cross the runtime boundary")
    report = {"source_hashes": inputs.used, "artifacts": artifacts, "versions": versions, "runtime": boundary,
              "design_sha256": digest((HERE / "presentation.json").read_bytes()),
              "implementation": {str(p.relative_to(ROOT)): digest(p.read_bytes()) for p in [Path(__file__), *sorted((HERE / "ui").glob("workbench.*"))]},
              "model_sha256": digest(canonical(payload))}
    (output / "verification.json").write_bytes(canonical(report))
    return report


def record(output: Path, report: dict[str, Any]) -> None:
    from proofs import chains
    sys.path.insert(0, str(HERE))
    import itil

    if output.resolve() != DEFAULT_OUT.resolve():
        raise ValueError("Only the canonical output directory can be registered")
    itil.load_register()
    chains.PROOFS.clear()
    for a in report["artifacts"]:
        p = output / a["file"]
        state = itil.check_register(str(p.relative_to(ROOT)), a["sha256"])
        pr = chains.Proof(a["pid"], a["title"], [", ".join(a["sources"])],
                          [chains.Step("verify inputs", "hashlib.sha256 against proofs/register.json", lambda c: None),
                           chains.Step("transform and render", a["functions"], lambda c: None),
                           chains.Step("read back and compare", "lxml.html.fromstring; json.loads; typed model equality", lambda c: None)],
                          "standalone interface")
        pr.evidence = [(a["label"], [(p, a["sha256"], state)], a["shows"] + ["embedded model roundtrip True; en-US; external assets 0; connect-src none"])]
        chains.PROOFS.append(pr)
    evidence = output / "verification.json"
    h = digest(evidence.read_bytes())
    state = itil.check_register(str(evidence.relative_to(ROOT)), h)
    p = chains.Proof("P71", "presentation chain -> guarded execution -> source and output verification", ["P65-P70; registered input artifacts; presentation.json"],
                     [chains.Step("guard", "importlib.abc.MetaPathFinder; sys.addaudithook", lambda c: None),
                      chains.Step("check", "hashlib.sha256; source register equality; lxml and JSON roundtrip", lambda c: None),
                      chains.Step("record versions", "importlib.metadata.version; json.dumps", lambda c: None)], "verification record")
    p.evidence = [("all", [(evidence, h, state)], [f"registered inputs checked {len(report['source_hashes'])}", "model import attempts 0; network and child-execution attempts 0", "version pins: " + ", ".join(f"{k} {v}" for k, v in report["versions"].items())])]
    model_path = output / "model.json"
    model_hash = digest(model_path.read_bytes())
    model_state = itil.check_register(str(model_path.relative_to(ROOT)), model_hash)
    p.evidence.append(("model", [(model_path, model_hash, model_state)], ["typed source model registered as an explicit input for later augmentation chains", "all six embedded view payloads checked against the same typed values"]))
    chains.PROOFS.append(p)
    preface = "## Presentation execution contract"
    if preface not in chains.MD.read_text():
        original = chains.MD.read_text()
        contract = "\n" + preface + "\n\n**Zero LLM reasoning intervention between functions or chains is the framework's core execution claim and requirement.** LLM reasoning is confined to upstream ideation. Every runtime handoff, decision, transformation, and presentation is executed by defined functions and data.\n\nThe initial augmentations prioritize front end usability and presentation of the core deliverables, then expand into additional capabilities. P65 onward implements this sequence. The presentation runner checks registered input bytes and rejects model-client imports and network calls. New interface text uses US English; quoted source text remains verbatim.\n\n[Open the combined workbench](proofs/out/augmentation/workbench.html). Validation details and current limits accompany each added proof.\n"
        chains.MD.write_text(original.replace("# Proofs\n", "# Proofs\n" + contract, 1))
    counts = chains.append_ledger()
    # Amendments have been appended above; advance only these checked outputs.
    itil.save_register(reregister=True)
    print(json.dumps({"ledger": counts, "new": len(itil.NEW), "changed": len(itil.CHANGED)}))


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUT)
    parser.add_argument("--record", action="store_true")
    args = parser.parse_args()
    boundary = install_boundary()
    report = build(args.output, boundary)
    if args.record:
        record(args.output, report)
    print(json.dumps({"artifacts": len(report["artifacts"]), "source_files": len(report["source_hashes"]), "model_sha256": report["model_sha256"]}))


if __name__ == "__main__":
    main()
