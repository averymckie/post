"""The proofs, as code.

Every proof is a chain of steps.  A step is a plain-English change bound to the
functions that perform it.  Running this file executes every chain on the real
inputs under packs/ and proofs/in/, writes each deliverable under proofs/out/,
and writes PROOFS.md at the repository root: the chains in arrow form, each
followed by its deliverables and their SHA-256 digests.

    .venv/bin/python proofs/chains.py
"""
from __future__ import annotations

import collections
import csv
import datetime as dt
import hashlib
import io
import json
import os
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "src"))
OUT = ROOT / "proofs" / "out"
MD = ROOT / "PROOFS.md"
os.environ.setdefault("PM4PY_SHOW_PROGRESS_BAR", "false")

# ----------------------------------------------------------------------------
# framework
# ----------------------------------------------------------------------------


@dataclass
class Step:
    change: str
    fns: str
    run: Callable[[dict[str, Any]], None]


@dataclass
class Proof:
    pid: str
    title: str
    inputs: list[str]  # lines naming what enters; two names on one line is a join
    steps: list[Step]
    result: str
    evidence: list[tuple[str, list[tuple[Path, str]], list[str]]] = field(default_factory=list)


PROOFS: list[Proof] = []
RESULTS: dict[str, Any] = {}


def sha256_file(p: Path) -> str:
    return hashlib.sha256(p.read_bytes()).hexdigest()


def sha256_json(obj: Any) -> str:
    return hashlib.sha256(json.dumps(obj, sort_keys=True, ensure_ascii=False, default=str).encode()).hexdigest()


def out_dir(pid: str, label: str) -> Path:
    d = OUT / pid / label
    d.mkdir(parents=True, exist_ok=True)
    return d


def write_json(p: Path, obj: Any) -> Path:
    p.write_text(json.dumps(obj, indent=1, sort_keys=True, ensure_ascii=False, default=str) + "\n", encoding="utf-8")
    return p


def write_jsonl(p: Path, rows: list[Any]) -> Path:
    with p.open("w", encoding="utf-8", newline="\n") as fh:
        for r in rows:
            fh.write(json.dumps(r, sort_keys=True, ensure_ascii=False, default=str) + "\n")
    return p


def workbook(p: Path, sheets: dict[str, list[list[Any]]]) -> Path:
    import openpyxl

    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    for name, rows in sheets.items():
        ws = wb.create_sheet(name[:31])
        for r in rows:
            ws.append([("" if v is None else v) if isinstance(v, (int, float, str)) or v is None else str(v) for v in r])
    wb.save(p)
    return p


def run(proof: Proof, label: str, ctx: dict[str, Any]) -> dict[str, Any]:
    for s in proof.steps:
        s.run(ctx)
    files = [(p, sha256_file(p)) for p in ctx.get("files", [])]
    existing = next((p for p in PROOFS if p.pid == proof.pid), None)
    if existing is None:
        PROOFS.append(proof)
        existing = proof
    existing.evidence.append((label, files, ctx.get("shows", [])))
    return ctx


STOP = set(
    "the a an of to in on at by for and or is are be was were been being with as that this these those it its "
    "from into than then there their they them he she his her we our you your i me my will would shall should "
    "may might must can could not no nor if but so such any all each other which who whom what when where how "
    "also more most very only same do does did done has have had having about over under between within after "
    "before upon per via etc one two three four five six seven eight nine ten".split()
)


def lit(s: Any) -> str:
    """A clingo string literal."""
    return '"' + str(s).replace("\\", "\\\\").replace('"', '\\"') + '"'


def words(text: str) -> set[str]:
    return {w for w in re.findall(r"[a-z][a-z\-']+", text.lower()) if w not in STOP and len(w) > 2}


REF_ATTRS = ("sourceRef", "targetRef", "source", "target", "bpmnElement", "idref", "processRef")


def canonical_xml(path: Path) -> None:
    """Rewrite the ids of a BPMN or PNML file so that structurally equal exports are byte-equal.

    pm4py names nodes with fresh uuids on every export.  Ids are recomputed here from structure
    alone (tag, name, markings, and the neighbourhood of references, refined four times), elements
    are ordered by those ids, and the file is written back.  Nothing semantic changes.
    """
    import lxml.etree as ET

    tree = ET.parse(str(path))
    root = tree.getroot()

    def local(e: Any) -> str:
        return e.tag.split("}")[-1] if isinstance(e.tag, str) else ""

    elems = [e for e in root.iter() if isinstance(e.tag, str) and e.get("id")]
    by_id = {e.get("id"): e for e in elems}

    def label(e: Any) -> str:
        n = e.get("name") or ""
        for c in e:
            if local(c) == "name":
                for t in c:
                    if local(t) == "text" and t.text:
                        n = t.text
            if local(c) in ("initialMarking",) and len(c) and c[0].text:
                n += f"|init={c[0].text.strip()}"
        return n

    def refs(e: Any) -> list[tuple[str, str]]:
        out = []
        for a in REF_ATTRS:
            v = e.get(a)
            if v in by_id:
                out.append((a, v))
        for c in e:
            if local(c) in ("incoming", "outgoing") and c.text and c.text.strip() in by_id:
                out.append((local(c), c.text.strip()))
            for a in REF_ATTRS:
                v = c.get(a)
                if v in by_id:
                    out.append((f"{local(c)}.{a}", v))
        return out

    links = {e.get("id"): refs(e) for e in elems}
    back: dict[str, list[tuple[str, str]]] = collections.defaultdict(list)
    for src, lst in links.items():
        for role, dst in lst:
            back[dst].append((role, src))
    sig = {i: hashlib.sha256(f"{local(e)}|{label(e)}".encode()).hexdigest() for i, e in by_id.items()}
    for _ in range(4):
        sig = {
            i: hashlib.sha256(json.dumps([sig[i], sorted((r, sig[d]) for r, d in links[i]), sorted((r, sig[s]) for r, s in back[i])]).encode()).hexdigest()
            for i in by_id
        }
    order = sorted(by_id, key=lambda i: (local(by_id[i]), sig[i]))
    counters: dict[str, int] = collections.defaultdict(int)
    new: dict[str, str] = {}
    for i in order:
        tag = local(by_id[i])
        counters[tag] += 1
        new[i] = f"{tag}_{counters[tag]}"
    for e in root.iter():
        if not isinstance(e.tag, str):
            continue
        if e.get("id") in new:
            e.set("id", new[e.get("id")])
        for a in REF_ATTRS:
            if e.get(a) in new:
                e.set(a, new[e.get(a)])
        if local(e) in ("incoming", "outgoing") and e.text and e.text.strip() in new:
            e.text = new[e.text.strip()]
    for parent in root.iter():
        if not isinstance(parent.tag, str):
            continue
        kids = list(parent)
        if len(kids) > 1 and all(isinstance(k.tag, str) for k in kids):
            keyed = [((local(k), k.get("id") or k.get("bpmnElement") or k.get("idref") or (k.text or "").strip()), k) for k in kids]
            if len({k[0] for k in keyed}) == len(keyed):
                for k in kids:
                    parent.remove(k)
                for _, k in sorted(keyed, key=lambda kv: kv[0]):
                    parent.append(k)
    tree.write(str(path), xml_declaration=True, encoding="UTF-8", pretty_print=True)


# ----------------------------------------------------------------------------
# door one: procedures -> facts ; facts -> ordered steps -> digest
# ----------------------------------------------------------------------------

from compiled_ai.pack import model_path, read_sources  # noqa: E402
from compiled_ai.parse import parse_source  # noqa: E402
from compiled_ai.fol import compile_atoms  # noqa: E402
from compiled_ai.normalize import normalize  # noqa: E402
from compiled_ai.check import check_atoms  # noqa: E402
from compiled_ai.adjudicate import apply_rejections, load_rejections  # noqa: E402
from compiled_ai.reconcile import reconcile  # noqa: E402
from compiled_ai.order import order as order_atoms  # noqa: E402
from compiled_ai.seal import manifest_json, seal  # noqa: E402

PACKS = {
    "usc5-552-doj": ROOT / "packs" / "foia",
    "nodejs-tsc-charter": ROOT / "packs" / "nodejs-tsc-charter",
    "nodejs-governance": ROOT / "packs" / "nodejs-governance",
    "nodejs-tsc-minutes": ROOT / "packs" / "nodejs-tsc-minutes",
}


def s_read(ctx: dict[str, Any]) -> None:
    sources, errs = read_sources(ctx["pack"])
    assert not errs, errs
    ctx["sources"] = sources


CACHE = ROOT / "proofs" / "cache"


def s_parse(ctx: dict[str, Any]) -> None:
    """Dependency parse, memoized on (source sha256, model sha256, canon version): the same bytes parse the same way."""
    from compiled_ai.model import ParsedSentence

    mp = model_path(ctx["pack"])
    model_sha = sha256_file(mp)
    parsed: list[Any] = []
    for s in ctx["sources"]:
        key = hashlib.sha256(f"{s.sha256}|{model_sha}|{s.canon_version}|{s.kind}|{s.id}".encode()).hexdigest()
        cache = CACHE / "parse" / f"{key}.json"
        if cache.exists():
            parsed.extend(ParsedSentence.model_validate(d) for d in json.loads(cache.read_text(encoding="utf-8")))
            continue
        ps = parse_source(s, mp)
        cache.parent.mkdir(parents=True, exist_ok=True)
        cache.write_text(json.dumps([x.model_dump() for x in ps], ensure_ascii=False), encoding="utf-8")
        parsed.extend(ps)
    ctx["parsed"] = parsed


def s_extract(ctx: dict[str, Any]) -> None:
    ctx["atoms"] = compile_atoms(ctx["parsed"])


def s_route(ctx: dict[str, Any]) -> None:
    ctx["atoms"], ctx["candidates"] = normalize(ctx["atoms"], ctx["parsed"])


def s_check(ctx: dict[str, Any]) -> None:
    check_atoms(ctx["atoms"], ctx["parsed"])


def s_adjudicate(ctx: dict[str, Any]) -> None:
    ctx["atoms"] = apply_rejections(ctx["atoms"], load_rejections(ctx["pack"] / "rejections.yaml"))
    d = ctx["dir"]
    p = write_jsonl(d / "facts.jsonl", [a.model_dump() for a in ctx["atoms"]])
    ctx["files"] = [p]
    ctx["shows"] = [f"facts {len(ctx['atoms'])}; sentences {len(ctx['parsed'])}"] + [
        f'{a.predicate}({", ".join(x.split("#")[-1] for x in a.args)})  "{a.quote}"' for a in ctx["atoms"][:4]
    ]


def door_one(label: str) -> dict[str, Any]:
    proof = Proof(
        "P1",
        "procedures -> facts",
        ["procedures"],
        [
            Step("read", "compiled_ai.read.read_html | compiled_ai.read.read_text; html.parser.HTMLParser", s_read),
            Step("canonical text", "compiled_ai.canon.canonicalize; unicodedata.normalize NFC", lambda c: None),
            Step("dependency parse", "compiled_ai.parse.parse_source; ufal.udpipe.Pipeline.process", s_parse),
            Step("predicate extraction", "compiled_ai.fol.compile_atoms; predpatt.PredPatt", s_extract),
            Step("project", "compiled_ai.fol.compile_atoms; projection table as data", lambda c: None),
            Step("route roles", "compiled_ai.normalize.normalize; clorm.clingo.Control, two rules as data", s_route),
            Step("byte check", "compiled_ai.check.check_atoms; sentence.text[lo:hi] == quote", s_check),
            Step("adjudicate", "compiled_ai.adjudicate.apply_rejections; yaml.safe_load", s_adjudicate),
        ],
        "facts",
    )
    ctx = {"pack": PACKS[label], "dir": out_dir("P1", label), "label": label}
    run(proof, label, ctx)
    RESULTS[("facts", label)] = ctx
    return ctx


def s_reconcile(ctx: dict[str, Any]) -> None:
    ctx["rec"] = reconcile(ctx["atoms"])


def s_order(ctx: dict[str, Any]) -> None:
    ctx["ordering"] = order_atoms(ctx["atoms"])
    d = ctx["dir"]
    o = ctx["ordering"]
    ctx["files"] = [write_json(d / "ordered_steps.json", o.model_dump())]
    lemma = {a.args[0]: a.args[1] for a in ctx["atoms"] if a.predicate == "event"}
    ctx["shows"] = [f"consistent: {ctx['rec'].consistent}; cycle: {list(o.cycle)}"] + [
        f"{lemma.get(a, '?')} -> {lemma.get(b, '?')}   forced" for a, b in o.forced[:4]
    ]


def s_seal(ctx: dict[str, Any]) -> None:
    m = seal(ctx["pack"].name, ctx["sources"], ctx["parsed"], ctx["atoms"], ctx["rec"], ctx["ordering"], model_path(ctx["pack"]))
    ctx["manifest"] = m
    p = ctx["dir"] / "manifest.json"
    p.write_text(manifest_json(m), encoding="utf-8")
    ctx["files"].append(p)
    ctx["shows"].append(f"digest {m.digest}")


def door_one_order(label: str) -> dict[str, Any]:
    src = RESULTS[("facts", label)]
    proof = Proof(
        "P2",
        "facts -> ordered steps -> digest",
        ["facts"],
        [
            Step("consistency proof", "compiled_ai.reconcile.reconcile; z3.Solver.check, z3.Solver.unsat_core", s_reconcile),
            Step("forced order", "compiled_ai.order.order; networkx.transitive_reduction, networkx.lexicographical_topological_sort", s_order),
            Step("seal", "compiled_ai.seal.seal; json.dumps, hashlib.sha256", s_seal),
        ],
        "ordered steps, digest",
    )
    ctx = dict(src)
    ctx["dir"] = out_dir("P2", label)
    ctx["files"] = []
    run(proof, label, ctx)
    RESULTS[("ordered steps", label)] = ctx
    return ctx


# ----------------------------------------------------------------------------
# P3 facts -> anchor dates
# ----------------------------------------------------------------------------


def s_timexy(ctx: dict[str, Any]) -> None:
    import spacy
    import timexy  # noqa: F401

    nlp = spacy.blank("en")
    nlp.add_pipe("timexy", config={"kb_id_type": "timex3", "label": "timexy", "overwrite": False})
    sents = {ps.sentence.id: ps.sentence.text for ps in ctx["parsed"]}
    rows: list[dict[str, Any]] = []
    for sid, text in sorted(sents.items()):
        for e in nlp(text).ents:
            rows.append({"sentence_id": sid, "span": e.text, "timex3": e.kb_id_, "sentence": text})
    ctx["anchors"] = rows
    p = write_jsonl(ctx["dir"] / "anchor_dates.jsonl", rows)
    ctx["files"] = [p]
    ctx["shows"] = [f'"{r["span"]}" -> {r["timex3"]}' for r in rows[:4]]


def anchor_dates(label: str) -> dict[str, Any]:
    src = RESULTS[("facts", label)]
    proof = Proof(
        "P3",
        "facts -> anchor dates",
        ["facts"],
        [Step("anchor dates", "timexy (spacy.blank('en').add_pipe('timexy')); TIMEX3 durations and dates", s_timexy)],
        "anchor dates",
    )
    ctx = dict(src)
    ctx["dir"] = out_dir("P3", label)
    run(proof, label, ctx)
    RESULTS[("anchor dates", label)] = ctx
    return ctx


# ----------------------------------------------------------------------------
# P4 facts -> required actions -> decision table -> policy
# ----------------------------------------------------------------------------

JDM_TEMPLATE = """{
 "contentType": "application/vnd.gorules.decision",
 "nodes": [
  {"id": "in", "type": "inputNode", "name": "request", "position": {"x": 0, "y": 0}},
  {"id": "dt", "type": "decisionTableNode", "name": "required actions", "position": {"x": 0, "y": 0},
   "content": {"hitPolicy": "first",
    "inputs": [{"id": "i1", "name": "action", "field": "action"}],
    "outputs": [{"id": "o1", "name": "required", "field": "required"}, {"id": "o2", "name": "rule", "field": "rule"}],
    "rules": [
{% for r in rows %}     {"_id": "r{{ loop.index }}", "_description": {{ r.quote|tojson }}, "i1": {{ (r.action|tojson)|tojson }}, "o1": "true", "o2": "'r{{ loop.index }}'"}{{ "," if not loop.last }}
{% endfor %}    ]}},
  {"id": "out", "type": "outputNode", "name": "decision", "position": {"x": 0, "y": 0}}
 ],
 "edges": [{"id": "e1", "sourceId": "in", "targetId": "dt", "type": "edge"}, {"id": "e2", "sourceId": "dt", "targetId": "out", "type": "edge"}]
}
"""


def s_required(ctx: dict[str, Any]) -> None:
    atoms = ctx["atoms"]
    lemma = {a.args[0]: a.args[1] for a in atoms if a.predicate == "event"}
    quote = {a.sentence_id: a.quote for a in atoms}
    rows = []
    for a in atoms:
        if a.predicate == "obligatory":
            e = a.args[0]
            rows.append({"event": e, "action": lemma.get(e, "?"), "sentence_id": a.sentence_id, "quote": a.quote})
    ctx["required"] = rows
    ctx["shows"] = [f'{r["action"]}  "{r["quote"]}"' for r in rows[:3]]


def s_decision_table(ctx: dict[str, Any]) -> None:
    import jinja2

    rows = ctx["required"]
    seen: set[str] = set()
    uniq = [r for r in rows if not (r["action"] in seen or seen.add(r["action"]))]
    jdm = jinja2.Environment().from_string(JDM_TEMPLATE).render(rows=uniq)
    json.loads(jdm)
    p = ctx["dir"] / "policy.jdm.json"
    p.write_text(jdm, encoding="utf-8")
    ctx["jdm"] = jdm
    ctx["files"] = [p]


def s_policy(ctx: dict[str, Any]) -> None:
    import zen

    d = zen.ZenEngine().create_decision(ctx["jdm"])
    sample = ctx["required"][0]["action"]
    r = d.evaluate({"action": sample})["result"]
    ctx["decision"] = d
    ctx["shows"].append(f"zen loads the table; evaluate({{action: {sample!r}}}) -> {r}")


def policy(label: str) -> dict[str, Any]:
    src = RESULTS[("facts", label)]
    proof = Proof(
        "P4",
        "facts -> required actions -> decision table -> policy",
        ["facts"],
        [
            Step("required actions", "the obligatory facts with their event lemma and quote", s_required),
            Step("decision table", "jinja2.Environment.from_string(...).render -> GoRules JDM", s_decision_table),
            Step("policy", "zen.ZenEngine.create_decision; zen.ZenDecision.evaluate", s_policy),
        ],
        "policy",
    )
    ctx = dict(src)
    ctx["dir"] = out_dir("P4", label)
    run(proof, label, ctx)
    RESULTS[("policy", label)] = ctx
    return ctx


# ----------------------------------------------------------------------------
# P5 ordered steps -> process model -> process
# ----------------------------------------------------------------------------


def s_bpmn_from_order(ctx: dict[str, Any]) -> None:
    from pm4py.objects.bpmn.exporter import exporter as bpmn_exporter
    from pm4py.objects.bpmn.obj import BPMN

    o = ctx["ordering"]
    lemma = {a.args[0]: a.args[1] for a in ctx["atoms"] if a.predicate == "event"}
    b = BPMN()
    nodes: dict[str, Any] = {}
    start, end = BPMN.StartEvent(name="start"), BPMN.EndEvent(name="end")
    b.add_node(start)
    b.add_node(end)
    inv = {e for e in o.order}
    preds = collections.defaultdict(set)
    succs = collections.defaultdict(set)
    for a, c in o.forced:
        preds[c].add(a)
        succs[a].add(c)
    entry: dict[str, Any] = {}
    exit_: dict[str, Any] = {}
    gateways = 0
    for e in o.order:
        t = BPMN.Task(name=f"{lemma.get(e, '?')} [{e.split(':', 1)[1]}]")
        nodes[e] = t
        b.add_node(t)
        entry[e], exit_[e] = t, t
        if len(preds[e]) > 1:
            g = BPMN.ParallelGateway(name=f"join {e.split(':')[-1]}", gateway_direction=BPMN.Gateway.Direction.CONVERGING)
            b.add_node(g)
            b.add_flow(BPMN.SequenceFlow(g, t))
            entry[e] = g
            gateways += 1
        if len(succs[e]) > 1:
            g = BPMN.ParallelGateway(name=f"split {e.split(':')[-1]}", gateway_direction=BPMN.Gateway.Direction.DIVERGING)
            b.add_node(g)
            b.add_flow(BPMN.SequenceFlow(t, g))
            exit_[e] = g
            gateways += 1
    for a, c in o.forced:
        b.add_flow(BPMN.SequenceFlow(exit_[a], entry[c]))
    sources = [e for e in o.order if not preds[e]]
    sinks = [e for e in o.order if not succs[e]]
    first_node: Any = start
    if len(sources) > 1:
        first_node = BPMN.ParallelGateway(name="split start", gateway_direction=BPMN.Gateway.Direction.DIVERGING)
        b.add_node(first_node)
        b.add_flow(BPMN.SequenceFlow(start, first_node))
        gateways += 1
    last_node: Any = end
    if len(sinks) > 1:
        last_node = BPMN.ParallelGateway(name="join end", gateway_direction=BPMN.Gateway.Direction.CONVERGING)
        b.add_node(last_node)
        b.add_flow(BPMN.SequenceFlow(last_node, end))
        gateways += 1
    for e in sources:
        b.add_flow(BPMN.SequenceFlow(first_node, entry[e]))
    for e in sinks:
        b.add_flow(BPMN.SequenceFlow(exit_[e], last_node))
    p = ctx["dir"] / "process.bpmn"
    bpmn_exporter.apply(b, str(p))
    canonical_xml(p)
    ctx["bpmn"] = b
    ctx["files"] = [p]
    ctx["shows"] = [f"tasks {len(inv)}, parallel gateways {gateways}, flows {len(b.get_flows())}; first flow {lemma.get(o.forced[0][0])} -> {lemma.get(o.forced[0][1])}" if o.forced else "no forced edges"]


def process_from_order(label: str) -> dict[str, Any]:
    src = RESULTS[("ordered steps", label)]
    proof = Proof(
        "P5",
        "ordered steps -> process model -> process",
        ["ordered steps"],
        [Step("process model", "pm4py.objects.bpmn.obj.BPMN (StartEvent, Task, EndEvent, SequenceFlow); pm4py.objects.bpmn.exporter.exporter.apply", s_bpmn_from_order)],
        "process",
    )
    ctx = dict(src)
    ctx["dir"] = out_dir("P5", label)
    run(proof, label, ctx)
    RESULTS[("process", label)] = ctx
    return ctx


# ----------------------------------------------------------------------------
# door two: records -> facts
# ----------------------------------------------------------------------------


def s_read_rows_csv(ctx: dict[str, Any]) -> None:
    with open(ctx["path"], encoding="utf-8", newline="") as fh:
        ctx["rows"] = list(csv.DictReader(fh))
    ctx["shows"] = [f"columns {list(ctx['rows'][0].keys())[:6]}...", f"first row {dict(list(ctx['rows'][0].items())[:3])}"]


def s_read_rows_roster(ctx: dict[str, Any]) -> None:
    import lxml.html
    import markdown

    html = markdown.markdown(Path(ctx["path"]).read_text(encoding="utf-8"))
    doc = lxml.html.fromstring(html)
    rows = []
    for h in doc.iter("h4"):
        if h.text_content().strip() == "TSC voting members":
            ul = next(s for s in h.itersiblings() if s.tag == "ul")
            for li in ul.iter("li"):
                t = " ".join(li.text_content().split())
                m = re.match(r"(\S+) - (.+?)(?: <| \(|$)", t)
                if m:
                    rows.append({"handle": m.group(1), "name": m.group(2).strip()})
    ctx["rows"] = rows
    ctx["shows"] = [f"first rows {rows[:2]}"]


def s_type_rows(ctx: dict[str, Any]) -> None:
    import pydantic

    cols = list(ctx["rows"][0].keys())
    M = pydantic.create_model(ctx["model"], **{c: (str, ...) for c in cols})
    ctx["typed"] = [M(**r) for r in ctx["rows"]]
    ctx["shows"].append(f"pydantic model {ctx['model']} with fields {cols[:4]}...")


def s_assert_rows(ctx: dict[str, Any]) -> None:
    import clingo

    facts = []
    for r in ctx["typed"]:
        d = r.model_dump()
        vals = ",".join(lit(d[c]) for c in ctx["cols"])
        facts.append(f'{ctx["pred"]}({vals}).')
    prog = "\n".join(facts) + "\n"
    ctl = clingo.Control(["0"])
    ctl.add("base", [], prog)
    ctl.ground([("base", [])])
    n = 0
    for sym in ctl.symbolic_atoms:
        n += 1
    p = ctx["dir"] / "facts.lp"
    p.write_text(prog, encoding="utf-8")
    ctx["facts_lp"] = prog
    ctx["files"] = [p]
    ctx["shows"].append(f"clingo grounds {n} facts; first {facts[0][:120]}")


def door_two(label: str, path: Path, reader: Callable, model: str, pred: str, cols: list[str]) -> dict[str, Any]:
    proof = Proof(
        "P6",
        "records -> facts",
        ["records"],
        [
            Step("read rows", "csv.DictReader | markdown.markdown + lxml.html.fromstring", reader),
            Step("type", "pydantic.create_model", s_type_rows),
            Step("assert", "clingo.Control.add; rows as ground facts", s_assert_rows),
        ],
        "facts",
    )
    ctx = {"path": path, "dir": out_dir("P6", label), "model": model, "pred": pred, "cols": cols, "label": label}
    run(proof, label, ctx)
    RESULTS[("records", label)] = ctx
    return ctx


# ----------------------------------------------------------------------------
# door four: event log -> log -> process model -> process diagram ; replay
# ----------------------------------------------------------------------------


def s_read_xes(ctx: dict[str, Any]) -> None:
    import pm4py

    log = pm4py.read_xes(str(ctx["path"]))
    ctx["log"] = log
    ctx["shows"] = [
        f"events {len(log)}, cases {log['case:concept:name'].nunique()}, activities {log['concept:name'].nunique()}",
        f"span {log['time:timestamp'].min().date()} to {log['time:timestamp'].max().date()}",
    ]


def s_discover(ctx: dict[str, Any]) -> None:
    import pm4py

    log = ctx["log"]
    cases = sorted(log["case:concept:name"].unique())
    train = log[log["case:concept:name"].isin(cases[: len(cases) // 2])]
    net, im, fm = pm4py.discover_petri_net_inductive(train)
    ctx["net"] = (net, im, fm)
    ctx["train_cases"] = len(cases) // 2
    starts = pm4py.get_start_activities(train)
    ctx["shows"].append(f"discovered from the first half of cases by id; start activities {starts}")


def s_chart(ctx: dict[str, Any]) -> None:
    import pm4py
    from pm4py.objects.bpmn.exporter import exporter as bpmn_exporter
    from pm4py.objects.conversion.wf_net.variants import to_bpmn

    net, im, fm = ctx["net"]
    p1 = ctx["dir"] / "process.pnml"
    pm4py.write_pnml(net, im, fm, str(p1))
    canonical_xml(p1)
    b = to_bpmn.apply(net, im, fm)
    p2 = ctx["dir"] / "process.bpmn"
    bpmn_exporter.apply(b, str(p2))
    canonical_xml(p2)
    ctx["files"] = [p1, p2]
    ctx["shows"].append(f"petri net places {len(net.places)}, transitions {len(net.transitions)}; bpmn nodes {len(b.get_nodes())}")


def door_four(label: str, path: Path) -> dict[str, Any]:
    proof = Proof(
        "P8",
        "event log -> log -> process model -> process diagram",
        ["event log"],
        [
            Step("read", "pm4py.read_xes", s_read_xes),
            Step("discover", "pm4py.discover_petri_net_inductive", s_discover),
            Step("chart", "pm4py.write_pnml; pm4py.objects.conversion.wf_net.variants.to_bpmn.apply; pm4py.objects.bpmn.exporter.exporter.apply", s_chart),
        ],
        "process diagram",
    )
    ctx = {"path": path, "dir": out_dir("P8", label), "label": label}
    run(proof, label, ctx)
    RESULTS[("log", label)] = ctx
    return ctx


def s_replay(ctx: dict[str, Any]) -> None:
    import pm4py

    net, im, fm = ctx["net"]
    log = ctx["log"]
    el = pm4py.convert_to_event_log(log)
    res = pm4py.conformance_diagnostics_token_based_replay(el, net, im, fm)
    rows = [["case", "trace_fitness", "is_fit", "missing_tokens", "remaining_tokens"]]
    fit = 0
    for tr, r in zip(el, res):
        rows.append([tr.attributes["concept:name"], round(r["trace_fitness"], 4), r["trace_is_fit"], r["missing_tokens"], r["remaining_tokens"]])
        fit += bool(r["trace_is_fit"])
    p = workbook(ctx["dir"] / "conformance.xlsx", {"conformance": rows})
    ctx["conformance_rows"] = rows
    ctx["files"] = [p]
    ctx["shows"] = [f"cases fit {fit} of {len(res)}; model from the first {ctx['train_cases']} cases", f"first rows {rows[1:3]}"]


def replay(label: str) -> dict[str, Any]:
    src = RESULTS[("log", label)]
    proof = Proof(
        "P9",
        "process, event log -> conformance -> workbook",
        ["process, event log"],
        [
            Step("replay", "pm4py.convert_to_event_log; pm4py.conformance_diagnostics_token_based_replay", s_replay),
            Step("tabulate", "openpyxl.Workbook.save", lambda c: None),
        ],
        "conformance",
    )
    ctx = dict(src)
    ctx["dir"] = out_dir("P9", label)
    run(proof, label, ctx)
    RESULTS[("conformance", label)] = ctx
    return ctx


# ----------------------------------------------------------------------------
# door three: minutes -> parsed minutes (the same functions as door one)
# ----------------------------------------------------------------------------


def door_three(label: str) -> dict[str, Any]:
    proof = Proof(
        "P7",
        "minutes -> parsed minutes",
        ["minutes"],
        [
            Step("read", "compiled_ai.read.read_text", s_read),
            Step("dependency parse", "compiled_ai.parse.parse_source; ufal.udpipe.Pipeline.process", s_parse),
            Step("predicate extraction", "compiled_ai.fol.compile_atoms; predpatt.PredPatt", s_extract),
            Step("route roles", "compiled_ai.normalize.normalize; clorm.clingo.Control", s_route),
            Step("byte check", "compiled_ai.check.check_atoms", s_check),
            Step("adjudicate", "compiled_ai.adjudicate.apply_rejections", s_adjudicate),
        ],
        "parsed minutes",
    )
    ctx = {"pack": PACKS[label], "dir": out_dir("P7", label), "label": label}
    run(proof, label, ctx)
    RESULTS[("parsed minutes", label)] = ctx
    return ctx


# ----------------------------------------------------------------------------
# joins
# ----------------------------------------------------------------------------


def step_words(ctx: dict[str, Any], events: list[str]) -> dict[str, set[str]]:
    atoms = ctx["atoms"]
    ws: dict[str, set[str]] = {e: set() for e in events}
    for a in atoms:
        if a.predicate == "event" and a.args[0] in ws:
            ws[a.args[0]].add(a.args[1].lower())
        if a.predicate in ("agent", "patient", "theme") and a.args[0] in ws:
            ws[a.args[0]] |= words(a.quote)
    return ws


def s_tag(ctx: dict[str, Any]) -> None:
    import clingo

    steps = ctx["steps"]  # dict event -> words
    minutes = ctx["minutes_ctx"]
    sents = {ps.sentence.id: ps.sentence.text for ps in minutes["parsed"]}
    line_words: dict[str, set[str]] = collections.defaultdict(set)
    for a in minutes["atoms"]:
        if a.predicate == "event":
            line_words[a.sentence_id].add(a.args[1].lower())
        elif a.predicate in ("agent", "patient", "theme"):
            line_words[a.sentence_id] |= words(a.quote)
    prog = []
    for e, ws in steps.items():
        for w in ws:
            prog.append(f'step_word({lit(e)},{lit(w)}).')
    for sid, ws in line_words.items():
        for w in ws:
            prog.append(f'line_word({lit(sid)},{lit(w)}).')
    prog.append("shared(L,E,W) :- line_word(L,W), step_word(E,W).")
    prog.append("score(L,E,N) :- line_word(L,_), step_word(E,_), N = #count{ W : shared(L,E,W) }, N > 0.")
    prog.append("best(L,E) :- score(L,E,N), N = #max{ M : score(L,_,M) }.")
    prog.append("tie(L) :- best(L,E1), best(L,E2), E1 < E2.")
    prog.append("untagged(L) :- line_word(L,_), not score(L,_,_).")
    prog.append("#show best/2. #show shared/3. #show tie/1. #show untagged/1.")
    ctl = clingo.Control(["0"])
    ctl.add("base", [], "\n".join(prog))
    ctl.ground([("base", [])])
    syms: list[Any] = []
    ctl.solve(on_model=lambda m: syms.extend(m.symbols(shown=True)))
    best = collections.defaultdict(list)
    shared = collections.defaultdict(set)
    ties, untagged = set(), set()
    for s in syms:
        if s.name == "best":
            best[s.arguments[0].string].append(s.arguments[1].string)
        elif s.name == "shared":
            shared[(s.arguments[0].string, s.arguments[1].string)].add(s.arguments[2].string)
        elif s.name == "tie":
            ties.add(s.arguments[0].string)
        elif s.name == "untagged":
            untagged.add(s.arguments[0].string)
    lemma = ctx["lemma"]
    rows = [["sentence_id", "sentence", "step", "step_lemma", "shared_words", "tie"]]
    tags = []
    for sid in sorted(best):
        for e in sorted(best[sid]):
            rows.append([sid, sents[sid], e, lemma.get(e, "?"), " ".join(sorted(shared[(sid, e)])), sid in ties])
            tags.append({"sentence_id": sid, "step": e, "shared": sorted(shared[(sid, e)]), "tie": sid in ties})
    unt = [["sentence_id", "sentence"]] + [[sid, sents[sid]] for sid in sorted(untagged)]
    p = workbook(ctx["dir"] / "tagged_steps.xlsx", {"tagged": rows, "untagged": unt})
    write_jsonl(ctx["dir"] / "tags.jsonl", tags)
    ctx["tags"] = tags
    ctx["files"] = [p, ctx["dir"] / "tags.jsonl"]
    strongest = sorted(rows[1:], key=lambda r: (-len(r[4].split()), r[0]))[:4]
    ctx["shows"] = [
        f"tagged sentences {len(best)}, ties {len(ties)}, untagged {len(untagged)}",
        *[f'{r[3]} <- "{r[1][:70]}"  shared: {r[4]}' for r in strongest],
    ]


def tag(label_steps: str, label_minutes: str, pid: str = "P10", required: bool = False) -> dict[str, Any]:
    steps_ctx = RESULTS[("ordered steps", label_steps)]
    minutes = RESULTS[("parsed minutes", label_minutes)]
    if required:
        events = sorted({a.args[0] for a in steps_ctx["atoms"] if a.predicate == "obligatory"})
        title, inputs = "required actions, parsed minutes -> tagged actions -> workbook", ["required actions (P4), parsed minutes"]
    else:
        events = list(steps_ctx["ordering"].order)
        title, inputs = "ordered steps, parsed minutes -> tagged steps -> workbook", ["ordered steps, parsed minutes"]
    proof = Proof(
        pid,
        title,
        inputs,
        [
            Step("tag", "clingo: shared words, #max score, ties flagged, untagged listed; stopword list as data", s_tag),
            Step("tabulate", "openpyxl.Workbook.save", lambda c: None),
        ],
        "tagged actions" if required else "tagged steps",
    )
    ctx = {
        "dir": out_dir(pid, f"{label_steps}+{label_minutes}"),
        "steps": step_words(steps_ctx, events),
        "minutes_ctx": minutes,
        "lemma": {a.args[0]: a.args[1] for a in steps_ctx["atoms"] if a.predicate == "event"},
        "order": events,
        "forced": list(steps_ctx["ordering"].forced),
    }
    run(proof, f"{label_steps}+{label_minutes}", ctx)
    RESULTS[("tagged steps", label_steps)] = ctx
    return ctx


def s_measure_tags(ctx: dict[str, Any]) -> None:
    import clingo

    prog = [f'tag({lit(t["sentence_id"])},{lit(t["step"])}).' for t in ctx["tags"]]
    prog += [f'step({lit(e)}).' for e in ctx["order"]]
    prog.append("count(E,N) :- step(E), N = #count{ L : tag(L,E) }.")
    prog.append("#show count/2.")
    ctl = clingo.Control(["0"])
    ctl.add("base", [], "\n".join(prog))
    ctl.ground([("base", [])])
    counts: dict[str, int] = {}
    ctl.solve(on_model=lambda m: counts.update({s.arguments[0].string: s.arguments[1].number for s in m.symbols(shown=True)}))
    rows = [["step", "step_lemma", "minutes_sentences"]] + [[e, ctx["lemma"].get(e, "?"), counts.get(e, 0)] for e in ctx["order"]]
    p = workbook(ctx["dir"] / "measured_steps.xlsx", {"measured": rows})
    ctx["counts"] = counts
    ctx["files"] = [p]
    ctx["shows"] = [f"{r[1]}: {r[2]}" for r in sorted(rows[1:], key=lambda r: -r[2])[:4]]


def measure_tags(label_steps: str) -> dict[str, Any]:
    src = RESULTS[("tagged steps", label_steps)]
    proof = Proof(
        "P11",
        "tagged steps -> measured steps -> workbook",
        ["tagged steps"],
        [Step("measure", "clingo #count per step", s_measure_tags), Step("tabulate", "openpyxl.Workbook.save", lambda c: None)],
        "measured steps",
    )
    ctx = dict(src)
    ctx["dir"] = out_dir("P11", label_steps)
    run(proof, label_steps, ctx)
    RESULTS[("measured steps", label_steps)] = ctx
    return ctx


def s_key_measure_log(ctx: dict[str, Any]) -> None:
    import duckdb

    net, _, _ = ctx["net"]
    labels = sorted(t.label for t in net.transitions if t.label)
    rows = ctx["rows"]
    con = duckdb.connect()
    con.execute("create table ev(case_id varchar, activity varchar, ts timestamptz)")
    con.executemany("insert into ev values (?, ?, ?)", [(r["case:concept:name"], r["concept:name"], r["time:timestamp"]) for r in rows])
    con.execute("create table model(activity varchar)")
    con.executemany("insert into model values (?)", [(l,) for l in labels])
    res = con.execute(
        """
        with seq as (
          select case_id, activity, ts,
                 lag(ts) over (partition by case_id order by ts) as prev
          from ev)
        select m.activity, count(e.activity) as events,
               round(avg(epoch(e.ts - e.prev))/3600, 2) as mean_hours_since_previous
        from model m left join seq e on e.activity = m.activity
        group by m.activity order by m.activity
        """
    ).fetchall()
    table = [["activity", "events", "mean_hours_since_previous"]] + [list(r) for r in res]
    p = workbook(ctx["dir"] / "measured_activities.xlsx", {"measured": table})
    ctx["measured"] = table
    ctx["files"] = [p]
    ctx["shows"] = [f"{r[0]}: events {r[1]}, mean hours since previous {r[2]}" for r in res[:4]]


def key_measure_log(label_log: str, label_records: str) -> dict[str, Any]:
    proof = Proof(
        "P12",
        "process, records -> measured steps -> workbook",
        ["process, records"],
        [
            Step("key", "duckdb: join on activity name between the model's transition labels and the rows", s_key_measure_log),
            Step("measure", "duckdb: count, avg(epoch(ts - lag(ts)))", lambda c: None),
            Step("tabulate", "openpyxl.Workbook.save", lambda c: None),
        ],
        "measured steps",
    )
    ctx = {"dir": out_dir("P12", f"{label_log}+{label_records}"), "net": RESULTS[("log", label_log)]["net"], "rows": RESULTS[("records", label_records)]["rows"]}
    run(proof, f"{label_log}+{label_records}", ctx)
    RESULTS[("measured activities", f"{label_log}+{label_records}")] = ctx
    return ctx


MAJORITY_JDM = """{
 "contentType": "application/vnd.gorules.decision",
 "nodes": [
  {"id": "in", "type": "inputNode", "name": "meeting", "position": {"x": 0, "y": 0}},
  {"id": "dt", "type": "decisionTableNode", "name": "simple majority of all TSC voting members", "position": {"x": 0, "y": 0},
   "content": {"hitPolicy": "first",
    "inputs": [{"id": "i1", "name": "present voting members", "field": "present"}],
    "outputs": [{"id": "o1", "name": "majority reachable", "field": "majority_reachable"}, {"id": "o2", "name": "rule", "field": "rule"}],
    "rules": [
     {"_id": "r1", "_description": {{ quote|tojson }}, "i1": ">= {{ majority }}", "o1": "true", "o2": "'r1'"},
     {"_id": "r2", "_description": {{ quote|tojson }}, "i1": "", "o1": "false", "o2": "'r2'"}
    ]}},
  {"id": "out", "type": "outputNode", "name": "decision", "position": {"x": 0, "y": 0}}
 ],
 "edges": [{"id": "e1", "sourceId": "in", "targetId": "dt", "type": "edge"}, {"id": "e2", "sourceId": "dt", "targetId": "out", "type": "edge"}]
}
"""


def s_attendance(ctx: dict[str, Any]) -> None:
    import lxml.html
    import markdown

    rows = []
    for src in ctx["minutes_ctx"]["sources"]:
        html = markdown.markdown(Path(src.path).read_text(encoding="utf-8"))
        doc = lxml.html.fromstring(html)
        for h in doc.iter("h2"):
            if h.text_content().strip() == "Present":
                ul = next(s for s in h.itersiblings() if s.tag == "ul")
                present = [" ".join(li.text_content().split()) for li in ul.iter("li")]
                voting = [p for p in present if "(voting member)" in p]
                rows.append({"meeting": src.id, "present_voting": len(voting), "present": len(present)})
    ctx["attendance"] = rows
    ctx["shows"] = [f'{r["meeting"]}: voting members present {r["present_voting"]}' for r in rows[:3]]


def s_majority_policy(ctx: dict[str, Any]) -> None:
    import jinja2
    import zen

    total = len(ctx["roster"])
    majority = total // 2 + 1
    rule = next(a for a in ctx["charter_atoms"] if a.predicate == "obligatory" and "majority" in a.quote.lower()) if any(
        a.predicate == "obligatory" and "majority" in a.quote.lower() for a in ctx["charter_atoms"]
    ) else None
    sent = {ps.sentence.id: ps.sentence.text for ps in ctx["charter_parsed"]}
    quote = next(t for t in sent.values() if "simple majority of all TSC voting members" in t)
    jdm = jinja2.Environment().from_string(MAJORITY_JDM).render(majority=majority, quote=quote)
    p = ctx["dir"] / "policy.jdm.json"
    p.write_text(jdm, encoding="utf-8")
    ctx["decision"] = zen.ZenEngine().create_decision(jdm)
    ctx["majority"] = majority
    ctx["files"] = [p]
    ctx["shows"].append(f'roster: {total} voting members (nodejs/node README, retrieved 2026-09-04); majority = {majority}; rule quoted from the charter: "{quote[:90]}..."')


def s_evaluate(ctx: dict[str, Any]) -> None:
    rows = [["meeting", "present_voting", "majority_reachable"]]
    for r in ctx["attendance"]:
        res = ctx["decision"].evaluate({"present": r["present_voting"]})["result"]
        rows.append([r["meeting"], r["present_voting"], res["majority_reachable"]])
    p = workbook(ctx["dir"] / "decisions.xlsx", {"decisions": rows})
    ctx["decision_rows"] = rows
    ctx["files"].append(p)
    ctx["shows"] += [f"{r[0]}: present {r[1]} -> majority reachable {r[2]}" for r in rows[1:]]


def evaluate_majority() -> dict[str, Any]:
    proof = Proof(
        "P13",
        "policy, records -> decisions -> workbook",
        ["policy (charter: simple majority of all TSC voting members), records (roster; attendance from the minutes)"],
        [
            Step("attendance rows", "markdown.markdown; lxml.html.fromstring; the Present list of each minutes file", s_attendance),
            Step("policy", "jinja2 render of the rule with the roster count -> GoRules JDM; zen.ZenEngine.create_decision", s_majority_policy),
            Step("evaluate", "zen.ZenDecision.evaluate per meeting", s_evaluate),
            Step("tabulate", "openpyxl.Workbook.save", lambda c: None),
        ],
        "decisions",
    )
    ctx = {
        "dir": out_dir("P13", "charter+roster+minutes"),
        "minutes_ctx": RESULTS[("parsed minutes", "nodejs-tsc-minutes")],
        "roster": RESULTS[("records", "tsc-voting-members")]["rows"],
        "charter_atoms": RESULTS[("facts", "nodejs-tsc-charter")]["atoms"],
        "charter_parsed": RESULTS[("facts", "nodejs-tsc-charter")]["parsed"],
    }
    run(proof, "charter+roster+minutes", ctx)
    RESULTS[("decisions", "majority")] = ctx
    return ctx


# ----------------------------------------------------------------------------
# renders
# ----------------------------------------------------------------------------


def draw_deck(path: Path, order: list[str], forced: list[tuple[str, str]], label_of: Callable[[str], str], title: str) -> dict[str, Any]:
    import networkx as nx
    from pptx import Presentation
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.util import Inches, Pt

    G = nx.DiGraph()
    G.add_nodes_from(order)
    G.add_edges_from(forced)
    layers = list(nx.topological_generations(G))
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(16), Inches(9)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(15), Inches(0.5))
    tb.text_frame.text = title
    pos: dict[str, Any] = {}
    for li, layer in enumerate(layers):
        for ni, n in enumerate(sorted(layer)):
            sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3 + li * 2.4), Inches(0.8 + ni * 0.62), Inches(2.1), Inches(0.5))
            sh.text_frame.text = label_of(n)
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    r.font.size = Pt(8)
            pos[n] = sh
    for a, b in forced:
        c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 0, 0, 0, 0)
        c.begin_connect(pos[a], 3)
        c.end_connect(pos[b], 1)
    prs.save(path)
    return {"layers": len(layers), "boxes": len(pos), "connectors": len(forced)}


def s_deck(ctx: dict[str, Any]) -> None:
    o = ctx["ordering"]
    lemma = {a.args[0]: a.args[1] for a in ctx["atoms"] if a.predicate == "event"}
    p = ctx["dir"] / "dependencies.pptx"
    info = draw_deck(p, list(o.order), list(o.forced), lambda e: f"{lemma.get(e, '?')}\n{e.split(':')[-1]}", f"{ctx['label']}: forced order")
    ctx["files"] = [p]
    ctx["shows"] = [f"layers {info['layers']}, boxes {info['boxes']}, connectors {info['connectors']}"]


def deck(label: str) -> dict[str, Any]:
    src = RESULTS[("ordered steps", label)]
    proof = Proof(
        "P14",
        "ordered steps -> deck",
        ["ordered steps"],
        [
            Step("layers", "networkx.topological_generations", lambda c: None),
            Step("draw", "pptx.Presentation; Shapes.add_shape, Shapes.add_connector, Connector.begin_connect, Connector.end_connect; Presentation.save", s_deck),
        ],
        "deck",
    )
    ctx = dict(src)
    ctx["dir"] = out_dir("P14", label)
    run(proof, label, ctx)
    RESULTS[("deck", label)] = ctx
    return ctx


def s_document(ctx: dict[str, Any]) -> None:
    import docx

    o = ctx["ordering"]
    atoms = ctx["atoms"]
    lemma = {a.args[0]: a.args[1] for a in atoms if a.predicate == "event"}
    quote = {a.args[0]: a.quote for a in atoms if a.predicate == "event"}
    sent = {ps.sentence.id: ps.sentence.text for ps in ctx["parsed"]}
    d = docx.Document()
    d.add_heading(f"{ctx['label']}: ordered steps", 1)
    for i, e in enumerate(o.order, 1):
        d.add_paragraph(f"{i}. {lemma.get(e, '?')}  [{e}]")
        d.add_paragraph(sent.get(e.split('#')[0], ""), style="Intense Quote")
    d.add_heading("forced precedence", 2)
    t = d.add_table(rows=1, cols=2)
    t.rows[0].cells[0].text, t.rows[0].cells[1].text = "before", "after"
    for a, b in o.forced:
        r = t.add_row().cells
        r[0].text, r[1].text = lemma.get(a, "?"), lemma.get(b, "?")
    p = ctx["dir"] / "ordered_steps.docx"
    d.save(p)
    ctx["files"] = [p]
    ctx["shows"] = [f"paragraphs {len(d.paragraphs)}, table rows {len(t.rows)}"]


def document(label: str) -> dict[str, Any]:
    src = RESULTS[("ordered steps", label)]
    proof = Proof("P15", "ordered steps -> document", ["ordered steps"], [Step("write", "docx.Document; Document.add_heading, add_paragraph, add_table; docx.document.Document.save", s_document)], "document")
    ctx = dict(src)
    ctx["dir"] = out_dir("P15", label)
    run(proof, label, ctx)
    return ctx


def s_tabulate_facts(ctx: dict[str, Any]) -> None:
    rows = [["id", "predicate", "args", "sentence_id", "quote"]] + [[a.id, a.predicate, " ".join(a.args), a.sentence_id, a.quote] for a in ctx["atoms"]]
    p = workbook(ctx["dir"] / "facts.xlsx", {"facts": rows})
    ctx["files"] = [p]
    ctx["shows"] = [f"rows {len(rows) - 1}; first {rows[1][1]} {rows[1][2]} \"{rows[1][4]}\""]


def tabulate(label: str) -> dict[str, Any]:
    src = RESULTS[("facts", label)]
    proof = Proof("P16", "facts -> workbook", ["facts"], [Step("tabulate", "openpyxl.Workbook; Worksheet.append; Workbook.save", s_tabulate_facts)], "workbook")
    ctx = dict(src)
    ctx["dir"] = out_dir("P16", label)
    run(proof, label, ctx)
    return ctx


def s_seal_all(ctx: dict[str, Any]) -> None:
    entries = []
    for pr in PROOFS:
        for label, files, _ in pr.evidence:
            for p, h in files:
                entries.append({"proof": pr.pid, "input": label, "file": str(p.relative_to(ROOT)), "sha256": h})
    entries.sort(key=lambda e: e["file"])
    digest = sha256_json(entries)
    p = write_json(ctx["dir"] / "digests.json", {"entries": entries, "digest": digest})
    ctx["files"] = [p]
    ctx["shows"] = [f"deliverables {len(entries)}; digest {digest}"]


def seal_all() -> dict[str, Any]:
    proof = Proof("P17", "anything -> digest", ["every deliverable above"], [Step("seal", "json.dumps sort_keys; hashlib.sha256", s_seal_all)], "digest")
    ctx = {"dir": out_dir("P17", "all")}
    run(proof, "all", ctx)
    return ctx


# ----------------------------------------------------------------------------
# composition, reverse, minutes as a log
# ----------------------------------------------------------------------------


def s_select_meeting(ctx: dict[str, Any]) -> None:
    import clingo

    meeting = ctx["meeting"]
    prog = [f'tag({lit(t["sentence_id"])},{lit(t["step"])}).' for t in ctx["tags"]]
    prog += [f'step({lit(e)}).' for e in ctx["order"]]
    prog += [f'forced({lit(a)},{lit(b)}).' for a, b in ctx["forced"]]
    prog.append(f'meeting_line(L) :- tag(L,_), @startswith(L, "{meeting}") = 1.')
    prog.append("selected(E) :- tag(L,E), meeting_line(L).")
    prog.append("count(E,N) :- selected(E), N = #count{ L : tag(L,E), meeting_line(L) }.")
    prog.append("edge(A,B) :- forced(A,B), selected(A), selected(B).")
    prog.append("#show selected/1. #show count/2. #show edge/2.")

    class Ctx:
        def startswith(self, s: Any, p: Any) -> Any:
            return clingo.Number(1 if s.string.startswith(p.string) else 0)

    ctl = clingo.Control(["0"])
    ctl.add("base", [], "\n".join(prog))
    ctl.ground([("base", [])], context=Ctx())
    sel, cnt, edges = [], {}, []
    for_model: list[Any] = []
    ctl.solve(on_model=lambda m: for_model.extend(m.symbols(shown=True)))
    for s in for_model:
        if s.name == "selected":
            sel.append(s.arguments[0].string)
        elif s.name == "count":
            cnt[s.arguments[0].string] = s.arguments[1].number
        elif s.name == "edge":
            edges.append((s.arguments[0].string, s.arguments[1].string))
    ctx["selected"] = sorted(sel, key=ctx["order"].index)
    ctx["counts"] = cnt
    ctx["edges"] = sorted(edges)
    ctx["shows"] = [f"meeting {meeting}: steps discussed {len(sel)}, forced edges among them {len(edges)}"]


def s_compose_deck(ctx: dict[str, Any]) -> None:
    lemma = ctx["lemma"]
    p = ctx["dir"] / "deck.pptx"
    rows = [[e, lemma.get(e, "?"), ctx["counts"].get(e, 0)] for e in ctx["selected"]]
    info = draw_deck(p, ctx["selected"], ctx["edges"], lambda e: f"{lemma.get(e, '?')}  ({ctx['counts'].get(e, 0)} lines)\n{e.split(':')[-1]}", f"steps discussed at {ctx['meeting']}, with minutes lines per step")
    ctx["rows"] = rows
    ctx["files"] = [p]
    ctx["shows"].append(f"deck layers {info['layers']}, boxes {info['boxes']}, connectors {info['connectors']}")


def s_compose_seal(ctx: dict[str, Any]) -> None:
    d = sha256_json(ctx["rows"])
    p = write_json(ctx["dir"] / "digest.json", {"rows": ctx["rows"], "digest": d})
    ctx["digest"] = d
    ctx["files"].append(p)
    ctx["shows"].append(f"digest over the rows {d}")


def compose(label_steps: str, meeting: str) -> dict[str, Any]:
    src = RESULTS[("tagged steps", label_steps)]
    proof = Proof(
        "P18",
        "ordered steps, parsed minutes -> tagged steps -> measured steps -> selected steps -> deck -> digest",
        ["tagged steps (P10), forced order (P2)"],
        [
            Step("select", "clingo: steps tagged by lines of one meeting; forced edges among them; #count lines per step", s_select_meeting),
            Step("layers", "networkx.topological_generations", lambda c: None),
            Step("draw", "pptx.Presentation; add_shape, add_connector; Presentation.save", s_compose_deck),
            Step("seal", "json.dumps sort_keys; hashlib.sha256", s_compose_seal),
        ],
        "deck, digest",
    )
    ctx = dict(src)
    ctx["dir"] = out_dir("P18", f"{label_steps}+{meeting}")
    ctx["meeting"] = meeting
    run(proof, f"{label_steps}+{meeting}", ctx)
    RESULTS[("composed deck", label_steps)] = ctx
    return ctx


def s_read_shapes(ctx: dict[str, Any]) -> None:
    from pptx import Presentation

    prs = Presentation(ctx["deck_path"])
    rows = []
    for sh in prs.slides[0].shapes:
        if sh.__class__.__name__ == "Shape" and sh.has_text_frame and "(" in sh.text_frame.text:
            first, ident = sh.text_frame.text.split("\n")
            lemma, n = re.match(r"(.+?)  \((\d+) lines\)", first).groups()
            rows.append([ident, lemma, int(n)])
    ctx["rows_back"] = sorted(rows, key=lambda r: r[0])
    ctx["shows"] = [f"boxes read back {len(rows)}"]


def s_reverse_seal(ctx: dict[str, Any]) -> None:
    fwd = sorted([[e.split(":")[-1], l, n] for e, l, n in ctx["rows"]], key=lambda r: r[0])
    d = sha256_json(ctx["rows_back"])
    ctx["digest_back"] = d
    ctx["fwd_rows"] = fwd
    ctx["shows"].append(f"digest' {d}")


def s_compare(ctx: dict[str, Any]) -> None:
    from csv_diff import compare

    fwd = {r[0]: {"lemma": r[1], "lines": str(r[2])} for r in ctx["fwd_rows"]}
    back = {r[0]: {"lemma": r[1], "lines": str(r[2])} for r in ctx["rows_back"]}
    diff = compare(fwd, back)
    same = not (diff["added"] or diff["removed"] or diff["changed"])
    p = write_json(ctx["dir"] / "compare.json", {"digest_forward": sha256_json(fwd and ctx["fwd_rows"]), "digest_back": ctx["digest_back"], "diff": diff, "match": same})
    ctx["files"] = [p]
    RESULTS.setdefault(("roundtrips", "all"), []).append({"proof": "P19", "label": "deck", "match": bool(same), "changed": len(diff["changed"])})
    ctx["shows"].append(f"rows match: {same}; changed {len(diff['changed'])}, added {len(diff['added'])}, removed {len(diff['removed'])}")


def reverse(label_steps: str) -> dict[str, Any]:
    src = RESULTS[("composed deck", label_steps)]
    proof = Proof(
        "P19",
        "deck -> rows -> digest' ; digest', digest -> match",
        ["deck (P18)"],
        [
            Step("read shapes", "pptx.Presentation; slide.shapes, Shape.text_frame", s_read_shapes),
            Step("seal", "json.dumps sort_keys; hashlib.sha256", s_reverse_seal),
            Step("compare", "csv_diff.compare", s_compare),
        ],
        "match, or the differing cell",
    )
    ctx = dict(src)
    ctx["dir"] = out_dir("P19", label_steps)
    ctx["deck_path"] = src["files"][0]
    ctx["files"] = []
    run(proof, label_steps, ctx)
    return ctx


def s_minutes_log(ctx: dict[str, Any]) -> None:
    import clingo

    prog = [f'tag({lit(t["sentence_id"])},{lit(t["step"])}).' for t in ctx["tags"] if not t["tie"]]
    prog += [f'forced({lit(a)},{lit(b)}).' for a, b in ctx["forced"]]
    prog.append("meeting(L,M) :- tag(L,_), M = @meeting(L).")
    prog.append("occurs(M,E,L) :- tag(L,E), meeting(L,M).")
    prog.append("first(M,E,L) :- occurs(M,E,L), L = #min{ L2 : occurs(M,E,L2) }.")
    prog.append("violated(M,A,B,LA,LB) :- forced(A,B), first(M,A,LA), first(M,B,LB), LB < LA.")
    prog.append("respected(M,A,B) :- forced(A,B), first(M,A,LA), first(M,B,LB), LA < LB.")
    prog.append("#show violated/5. #show respected/3.")

    class Ctx:
        def meeting(self, l: Any) -> Any:
            return clingo.String(l.string.split(":")[0])

    ctl = clingo.Control(["0"])
    ctl.add("base", [], "\n".join(prog))
    ctl.ground([("base", [])], context=Ctx())
    syms: list[Any] = []
    ctl.solve(on_model=lambda m: syms.extend(m.symbols(shown=True)))
    lemma = ctx["lemma"]
    v = [["meeting", "before", "after", "line_of_before", "line_of_after"]]
    r = [["meeting", "before", "after"]]
    for s in syms:
        a = [x.string for x in s.arguments]
        if s.name == "violated":
            v.append([a[0], lemma.get(a[1], "?"), lemma.get(a[2], "?"), a[3], a[4]])
        else:
            r.append([a[0], lemma.get(a[1], "?"), lemma.get(a[2], "?")])
    p = workbook(ctx["dir"] / "conformance.xlsx", {"violated": v, "respected": r})
    ctx["files"] = [p]
    ctx["shows"] = [f"forced edges respected in the minutes {len(r) - 1}, violated {len(v) - 1}", *[f"{x[0]}: {x[1]} before {x[2]}" for x in r[1:3]], *[f"VIOLATED {x[0]}: {x[2]} discussed before {x[1]}" for x in v[1:3]]]


def minutes_as_log(label_steps: str) -> dict[str, Any]:
    src = RESULTS[("tagged steps", label_steps)]
    proof = Proof(
        "P20",
        "tagged steps -> event log ; process, event log -> conformance -> workbook",
        ["tagged steps (P10), forced order (P2)"],
        [
            Step("event log", "clingo: case = meeting, activity = tagged step (ties dropped), time = first sentence that mentions it", s_minutes_log),
            Step("replay", "clingo: violated(M,A,B) :- forced(A,B), first mention of B before first mention of A", lambda c: None),
            Step("tabulate", "openpyxl.Workbook.save", lambda c: None),
        ],
        "conformance",
    )
    ctx = dict(src)
    ctx["dir"] = out_dir("P20", label_steps)
    run(proof, label_steps, ctx)
    return ctx


# ----------------------------------------------------------------------------
# PROOFS.md
# ----------------------------------------------------------------------------

DOORS = [
    ("procedures", "packs/foia/sources/usdoj-foia.gov-foia-statute.html (5 U.S.C. 552, DOJ); packs/nodejs-tsc-charter/sources/tsc-TSC-Charter.md; packs/nodejs-governance/sources/GOVERNANCE.md"),
    ("records", "proofs/in/receipt.csv (WABO receipt phase, 8577 rows); proofs/in/node-README.md (TSC voting members)"),
    ("minutes", "packs/nodejs-tsc-minutes/sources/tsc-*.md (eight Node.js TSC meetings)"),
    ("event log", "proofs/in/receipt.xes (WABO receipt phase, 1434 cases)"),
]

PINS = "ufal.udpipe 1.4.0.1 · predpatt 1.0.1 · clingo 5.8.2 · clorm 1.6.3 · z3-solver 5.1.0.0 · networkx 3.6.1 · timexy 0.1.3 · spacy 3.8.16 · zen-engine 2.0.2 · Jinja2 3.1.6 · pm4py 2.7.23.8 · pydantic 2.13.5 · duckdb 1.5.5 · markdown 3.6 · lxml 6.1.3 · python-pptx 1.0.2 · python-docx 1.2.0 · openpyxl 3.1.2 · csv-diff 1.2 · PyYAML 6.0.3"


def write_md() -> None:
    lines = ["# Proofs", "", "```", "thing", "-> change (functions)", "-> thing", "a line naming two things is a join", "every function ran on the named input; every deliverable carries its sha256", "```", "", "## Doors", "", "```"]
    for name, where in DOORS:
        lines.append(f"{name:12s} {where}")
    lines += ["```", ""]
    if ITIL_REPORT.get("practices"):
        sys.path.insert(0, str(ROOT / "proofs"))
        import itil

        lines += itil.md_section(ITIL_REPORT)
    lines += ["## Variations: the numbered proofs", ""]
    for pr in sorted(PROOFS, key=lambda p: int(re.sub(r"\D", "", p.pid) or 0)):
        lines += [f"## {pr.pid}  {pr.title}", "", "```"]
        for inp in pr.inputs:
            lines.append(inp)
        for s in pr.steps:
            lines.append(f"-> {s.change} ({s.fns})")
        lines.append(f"-> {pr.result}")
        lines += ["```", ""]
        for label, files, shows in pr.evidence:
            lines.append(f"**{label}**")
            lines.append("")
            lines.append("```")
            for p, h in files:
                lines.append(f"{str(p.relative_to(ROOT)):58s} sha256 {h}")
            for sh in shows:
                lines.append(f"shows: {sh}")
            lines += ["```", ""]
    lines += ["## Pins", "", "```", PINS, "```", ""]
    MD.write_text("\n".join(lines), encoding="utf-8")


# ----------------------------------------------------------------------------
# P21 .. P30
# ----------------------------------------------------------------------------


def duck_events(rows: list[dict[str, Any]]) -> Any:
    import duckdb

    con = duckdb.connect()
    con.execute("create table ev(case_id varchar, activity varchar, ts timestamptz, deadline varchar, enddate varchar, department varchar, channel varchar)")
    con.executemany(
        "insert into ev values (?, ?, ?, ?, ?, ?, ?)",
        [(r["case:concept:name"], r["concept:name"], r["time:timestamp"], r["case:deadline"], r["case:enddate"], r["case:department"], r["case:channel"]) for r in rows],
    )
    return con


def s_deadlines(ctx: dict[str, Any]) -> None:
    con = duck_events(ctx["rows"])
    res = con.execute(
        """
        with c as (select case_id, any_value(deadline) as deadline, any_value(enddate) as enddate,
                          any_value(department) as department, any_value(channel) as channel, max(ts) as last_event
                   from ev group by case_id)
        select case_id, department, channel, deadline, enddate, last_event,
               case when enddate = '' then 'open'
                    when try_cast(enddate as timestamptz) > try_cast(deadline as timestamptz) then 'after deadline'
                    else 'by deadline' end as outcome
        from c order by case_id
        """
    ).fetchall()
    rows = [["case", "department", "channel", "deadline", "enddate", "last_event", "outcome"]] + [list(r) for r in res]
    counts = collections.Counter(r[6] for r in res)
    p = workbook(ctx["dir"] / "deadlines.xlsx", {"cases": rows})
    ctx["deadline_rows"] = rows
    ctx["files"] = [p]
    ctx["shows"] = [f"outcomes {dict(sorted(counts.items()))}", f"first rows {[r[:1] + r[6:] for r in res[:3]]}"]


def deadlines(label: str) -> dict[str, Any]:
    proof = Proof(
        "P21",
        "records -> measured cases -> workbook",
        ["records (receipt.csv: the rows carry case:deadline and case:enddate)"],
        [
            Step("key", "duckdb: group the rows by case", s_deadlines),
            Step("measure", "duckdb: enddate compared with deadline per case; open when enddate is empty", lambda c: None),
            Step("tabulate", "openpyxl.Workbook.save", lambda c: None),
        ],
        "measured cases",
    )
    ctx = {"dir": out_dir("P21", label), "rows": RESULTS[("records", label)]["rows"]}
    run(proof, label, ctx)
    RESULTS[("measured cases", label)] = ctx
    return ctx


def s_by_department(ctx: dict[str, Any]) -> None:
    con = duck_events(ctx["rows"])
    q = """
        with c as (select case_id, any_value({col}) as k, min(ts) as first, max(ts) as last, any_value(enddate) as enddate, any_value(deadline) as deadline
                   from ev group by case_id)
        select k, count(*) as cases, round(avg(epoch(last - first))/86400, 2) as mean_days_first_to_last,
               sum(case when enddate <> '' and try_cast(enddate as timestamptz) > try_cast(deadline as timestamptz) then 1 else 0 end) as after_deadline
        from c group by k order by k
        """
    sheets = {}
    for col in ("department", "channel"):
        res = con.execute(q.format(col=col)).fetchall()
        sheets[col] = [[col, "cases", "mean_days_first_to_last", "after_deadline"]] + [list(r) for r in res]
    p = workbook(ctx["dir"] / "by_department_and_channel.xlsx", sheets)
    ctx["files"] = [p]
    ctx["shows"] = [f"{r[0]}: cases {r[1]}, mean days {r[2]}, after deadline {r[3]}" for r in sheets["department"][1:4]] + [
        f"channel {r[0]}: cases {r[1]}, mean days {r[2]}" for r in sheets["channel"][1:3]
    ]


def by_department(label: str) -> dict[str, Any]:
    proof = Proof(
        "P22",
        "records -> measured groups -> workbook",
        ["records"],
        [
            Step("key", "duckdb: group by case, then by department and by channel", s_by_department),
            Step("measure", "duckdb: count, avg(epoch(last - first)), cases ended after deadline", lambda c: None),
            Step("tabulate", "openpyxl.Workbook.save", lambda c: None),
        ],
        "measured groups",
    )
    ctx = {"dir": out_dir("P22", label), "rows": RESULTS[("records", label)]["rows"]}
    run(proof, label, ctx)
    return ctx


def s_dfg(ctx: dict[str, Any]) -> None:
    import pm4py

    dfg, sa, ea = pm4py.discover_dfg(ctx["log"])
    edges = sorted(dfg.items(), key=lambda kv: (-kv[1], kv[0]))
    rows = [["from", "to", "count"]] + [[a, b, n] for (a, b), n in edges]
    starts = [["activity", "cases"]] + sorted(([a, n] for a, n in sa.items()), key=lambda r: -r[1])
    ends = [["activity", "cases"]] + sorted(([a, n] for a, n in ea.items()), key=lambda r: -r[1])
    p = workbook(ctx["dir"] / "directly_follows.xlsx", {"edges": rows, "start": starts, "end": ends})
    p2 = write_json(ctx["dir"] / "directly_follows.json", {"edges": [[a, b, n] for (a, b), n in edges], "start": sa, "end": ea})
    ctx["dfg_rows"] = rows
    ctx["files"] = [p, p2]
    ctx["shows"] = [f"edges {len(edges)}; strongest {rows[1][0]} -> {rows[1][1]} ({rows[1][2]})", f"end activities {ends[1:3]}"]


def dfg(label: str) -> dict[str, Any]:
    proof = Proof(
        "P23",
        "log -> directly-follows graph -> workbook",
        ["log (P8)"],
        [Step("discover", "pm4py.discover_dfg", s_dfg), Step("tabulate", "openpyxl.Workbook.save; json.dumps", lambda c: None)],
        "directly-follows graph",
    )
    ctx = dict(RESULTS[("log", label)])
    ctx["dir"] = out_dir("P23", label)
    run(proof, label, ctx)
    RESULTS[("dfg", label)] = ctx
    return ctx


def s_dashboard(ctx: dict[str, Any]) -> None:
    import plotly.express as px
    import plotly.io as pio

    table = ctx["measured"][1:]
    fig = px.bar(x=[r[0] for r in table], y=[r[1] for r in table], labels={"x": "activity", "y": "events"}, title="events per activity, receipt phase")
    fig.update_layout(xaxis_tickangle=-45)
    p = ctx["dir"] / "dashboard.html"
    pio.write_html(fig, str(p), include_plotlyjs="cdn", full_html=True, div_id="chart")
    ctx["files"] = [p]
    ctx["shows"] = [f"one bar per activity, {len(table)} bars; plotly.js from cdn"]


def dashboard(label: str) -> dict[str, Any]:
    proof = Proof(
        "P24",
        "measured steps -> chart -> page",
        ["measured steps (P12)"],
        [Step("chart", "plotly.express.bar; Figure.update_layout", s_dashboard), Step("page", "plotly.io.write_html(full_html, div_id)", lambda c: None)],
        "page",
    )
    ctx = {"dir": out_dir("P24", label), "measured": RESULTS[("measured activities", label)]["measured"]}
    run(proof, label, ctx)
    return ctx


def s_obligations_doc(ctx: dict[str, Any]) -> None:
    import docx

    sent = {ps.sentence.id: ps.sentence.text for ps in ctx["parsed"]}
    atoms = ctx["atoms"]
    agent = collections.defaultdict(list)
    for a in atoms:
        if a.predicate == "agent":
            agent[a.args[0]].append(a.quote)
    d = docx.Document()
    d.add_heading(f"{ctx['label']}: required actions", 1)
    t = d.add_table(rows=1, cols=3)
    t.rows[0].cells[0].text, t.rows[0].cells[1].text, t.rows[0].cells[2].text = "who", "must", "sentence"
    for r in ctx["required"]:
        c = t.add_row().cells
        c[0].text = "; ".join(agent.get(r["event"], []))
        c[1].text = r["action"]
        c[2].text = sent.get(r["sentence_id"], "")
    p = ctx["dir"] / "required_actions.docx"
    d.save(p)
    ctx["files"] = [p]
    ctx["shows"] = [f"table rows {len(t.rows) - 1}"] + [f'{"; ".join(agent.get(r["event"], []))} must {r["action"]}' for r in ctx["required"][:3]]


def obligations_document(label: str) -> dict[str, Any]:
    proof = Proof(
        "P25",
        "required actions -> document",
        ["required actions (P4), facts (P1)"],
        [Step("write", "docx.Document; Document.add_table; docx.document.Document.save", s_obligations_doc)],
        "document",
    )
    ctx = dict(RESULTS[("policy", label)])
    ctx["dir"] = out_dir("P25", label)
    run(proof, label, ctx)
    return ctx


def s_conf_by_department(ctx: dict[str, Any]) -> None:
    import duckdb

    con = duckdb.connect()
    con.execute("create table conf(case_id varchar, fitness double, is_fit boolean)")
    con.executemany("insert into conf values (?, ?, ?)", [(r[0], float(r[1]), bool(r[2])) for r in ctx["conformance_rows"][1:]])
    con.execute("create table attr(case_id varchar, department varchar, channel varchar)")
    seen = set()
    vals = []
    for r in ctx["rows"]:
        if r["case:concept:name"] not in seen:
            seen.add(r["case:concept:name"])
            vals.append((r["case:concept:name"], r["case:department"], r["case:channel"]))
    con.executemany("insert into attr values (?, ?, ?)", vals)
    sheets = {}
    for col in ("department", "channel"):
        res = con.execute(
            f"select a.{col}, count(*) as cases, sum(case when c.is_fit then 1 else 0 end) as fit, round(avg(c.fitness), 4) as mean_fitness "
            f"from conf c join attr a on a.case_id = c.case_id group by a.{col} order by a.{col}"
        ).fetchall()
        sheets[col] = [[col, "cases", "fit", "mean_fitness"]] + [list(r) for r in res]
    p = workbook(ctx["dir"] / "conformance_by_group.xlsx", sheets)
    ctx["files"] = [p]
    ctx["shows"] = [f"{r[0]}: cases {r[1]}, fit {r[2]}, mean fitness {r[3]}" for r in sheets["department"][1:4]]


def conformance_by_department(label_log: str, label_records: str) -> dict[str, Any]:
    proof = Proof(
        "P27",
        "conformance, records -> measured groups -> workbook",
        ["conformance (P9), records (P6)"],
        [
            Step("key", "duckdb: join on case id", s_conf_by_department),
            Step("measure", "duckdb: cases, fit cases, avg(fitness) per department and per channel", lambda c: None),
            Step("tabulate", "openpyxl.Workbook.save", lambda c: None),
        ],
        "measured groups",
    )
    ctx = {"dir": out_dir("P27", f"{label_log}+{label_records}"), "conformance_rows": RESULTS[("conformance", label_log)]["conformance_rows"], "rows": RESULTS[("records", label_records)]["rows"]}
    run(proof, f"{label_log}+{label_records}", ctx)
    return ctx


def s_decisions_discussion(ctx: dict[str, Any]) -> None:
    import clingo

    prog = [f"decision({lit(r[0])},{lit(str(r[2]).lower())})." for r in ctx["decision_rows"][1:]]
    prog += [f"tag({lit(t['sentence_id'])},{lit(t['step'])})." for t in ctx["tags"] if not t["tie"]]
    prog.append("meeting(L,M) :- tag(L,_), M = @meeting(L).")
    prog.append("discussed(M,E) :- tag(L,E), meeting(L,M).")
    prog.append("n(M,N) :- decision(M,_), N = #count{ E : discussed(M,E) }.")
    prog.append("#show decision/2. #show discussed/2. #show n/2.")

    class Ctx:
        def meeting(self, l: Any) -> Any:
            return clingo.String(l.string.split(":")[0])

    ctl = clingo.Control(["0"])
    ctl.add("base", [], "\n".join(prog))
    ctl.ground([("base", [])], context=Ctx())
    syms: list[Any] = []
    ctl.solve(on_model=lambda m: syms.extend(m.symbols(shown=True)))
    dec, disc, n = {}, collections.defaultdict(list), {}
    for s in syms:
        a = [x.string if x.type == clingo.SymbolType.String else x.number for x in s.arguments]
        if s.name == "decision":
            dec[a[0]] = a[1]
        elif s.name == "discussed":
            disc[a[0]].append(ctx["lemma"].get(a[1], "?"))
        else:
            n[a[0]] = a[1]
    rows = [["meeting", "majority_reachable", "required_actions_discussed", "actions"]] + [[m, dec[m], n.get(m, 0), ", ".join(sorted(disc.get(m, [])))] for m in sorted(dec)]
    p = workbook(ctx["dir"] / "decisions_with_discussion.xlsx", {"meetings": rows})
    ctx["files"] = [p]
    ctx["shows"] = [f"{r[0]}: majority {r[1]}, required actions discussed {r[2]}: {r[3][:60]}" for r in rows[1:5]]


def decisions_with_discussion() -> dict[str, Any]:
    proof = Proof(
        "P28",
        "decisions, tagged actions -> meetings -> workbook",
        ["decisions (P13), tagged actions (P26)"],
        [
            Step("key", "clingo: join on meeting id", s_decisions_discussion),
            Step("measure", "clingo #count of required actions discussed per meeting", lambda c: None),
            Step("tabulate", "openpyxl.Workbook.save", lambda c: None),
        ],
        "meetings",
    )
    tagged = RESULTS[("tagged steps", "nodejs-tsc-charter")]
    ctx = {"dir": out_dir("P28", "majority+charter-tags"), "decision_rows": RESULTS[("decisions", "majority")]["decision_rows"], "tags": tagged["tags"], "lemma": tagged["lemma"]}
    run(proof, "majority+charter-tags", ctx)
    return ctx


def s_mark_executable(ctx: dict[str, Any]) -> None:
    import lxml.etree as ET

    t = ET.parse(str(ctx["bpmn_path"]))
    ns = "{http://www.omg.org/spec/BPMN/20100524/MODEL}"
    for pr in t.getroot().iter(f"{ns}process"):
        pr.set("isExecutable", "true")
    for el in t.getroot().iter(f"{ns}task"):
        el.tag = f"{ns}manualTask"
    p = ctx["dir"] / "process.executable.bpmn"
    t.write(str(p), xml_declaration=True, encoding="UTF-8")
    canonical_xml(p)
    ctx["exec_path"] = p
    ctx["files"] = [p]


def s_execute(ctx: dict[str, Any]) -> None:
    from SpiffWorkflow.bpmn.parser.BpmnParser import BpmnParser
    from SpiffWorkflow.bpmn.workflow import BpmnWorkflow
    from SpiffWorkflow.util.task import TaskState

    parser = BpmnParser()
    parser.add_bpmn_file(str(ctx["exec_path"]))
    pid = list(parser.process_parsers)[0]
    wf = BpmnWorkflow(parser.get_spec(pid))
    wf.do_engine_steps()
    trace: list[str] = []
    n = 0
    while not wf.is_completed() and n < 5000:
        ready = sorted(wf.get_tasks(state=TaskState.READY), key=lambda t: t.task_spec.bpmn_name or t.task_spec.name)
        if not ready:
            break
        t0 = ready[0]
        trace.append(t0.task_spec.bpmn_name or t0.task_spec.name)
        t0.run()
        wf.do_engine_steps()
        n += 1
    if not trace:
        done = sorted(wf.get_tasks(state=TaskState.COMPLETED), key=lambda t: (getattr(t, "last_state_change", 0.0), t.task_spec.bpmn_name or t.task_spec.name))
        trace = [t.task_spec.bpmn_name or t.task_spec.name for t in done if (t.task_spec.bpmn_name or "").endswith("]")]
    ctx["trace"] = trace
    ctx["completed"] = wf.is_completed()
    first = {}
    for i, name in enumerate(trace):
        if "[" in name:
            first.setdefault(name.split(" [")[1].rstrip("]"), i)
    o = ctx["ordering"]
    ok = sum(1 for a, b in o.forced if first.get(a.split(":", 1)[1], 10**9) < first.get(b.split(":", 1)[1], -1))
    rows = [["position", "task"]] + [[i + 1, name] for i, name in enumerate(trace)]
    p = workbook(ctx["dir"] / "execution_trace.xlsx", {"trace": rows})
    ctx["files"].append(p)
    RESULTS.setdefault(("traces", "all"), {})[ctx["label"]] = list(trace)
    ctx["shows"] = [f"workflow completed: {ctx['completed']}; tasks run {len(trace)}", f"forced edges completed in order: {ok} of {len(o.forced)}", f"first tasks {trace[:4]}"]


def execute_process(label: str) -> dict[str, Any]:
    proof = Proof(
        "P29",
        "process -> executable process -> execution trace -> workbook",
        ["process (P5), forced order (P2)"],
        [
            Step("mark executable", "lxml.etree.parse; process.set('isExecutable', 'true'); each task renamed manualTask; ElementTree.write", s_mark_executable),
            Step("load", "SpiffWorkflow.bpmn.parser.BpmnParser.add_bpmn_file; BpmnParser.get_spec", s_execute),
            Step("execute", "SpiffWorkflow.bpmn.workflow.BpmnWorkflow; do_engine_steps; Task.run on ready tasks in name order", lambda c: None),
            Step("check", "every forced edge: first completion of the earlier task precedes the later", lambda c: None),
            Step("tabulate", "openpyxl.Workbook.save", lambda c: None),
        ],
        "execution trace",
    )
    src = RESULTS[("process", label)]
    ctx = {"dir": out_dir("P29", label), "bpmn_path": src["files"][0], "ordering": src["ordering"], "label": label}
    run(proof, label, ctx)
    return ctx


def s_roundtrip(ctx: dict[str, Any]) -> None:
    import python_calamine
    from csv_diff import compare

    rows = python_calamine.CalamineWorkbook.from_path(str(ctx["xlsx"])).get_sheet_by_name("facts").to_python()
    back = [{"id": r[0], "predicate": r[1], "args": str(r[2]).split(" "), "sentence_id": r[3], "quote": r[4]} for r in rows[1:]]
    fwd = [{"id": a.id, "predicate": a.predicate, "args": list(a.args), "sentence_id": a.sentence_id, "quote": a.quote} for a in ctx["atoms"]]
    d_f, d_b = sha256_json(fwd), sha256_json(back)
    diff = compare({r["id"]: {k: str(v) for k, v in r.items()} for r in fwd}, {r["id"]: {k: str(v) for k, v in r.items()} for r in back})
    same = d_f == d_b
    p = write_json(ctx["dir"] / "roundtrip.json", {"digest_forward": d_f, "digest_back": d_b, "match": same, "changed": diff["changed"][:5], "added": diff["added"][:5], "removed": diff["removed"][:5]})
    ctx["files"] = [p]
    RESULTS.setdefault(("roundtrips", "all"), []).append({"proof": "P30", "label": "workbook", "match": bool(same), "changed": len(diff["changed"])})
    ctx["shows"] = [f"rows read back {len(back)}; digest forward {d_f[:16]}…, back {d_b[:16]}…; match {same}; changed {len(diff['changed'])}"]


def workbook_roundtrip(label: str) -> dict[str, Any]:
    proof = Proof(
        "P30",
        "workbook -> facts' -> digest' ; digest', digest -> match",
        ["workbook (P16), facts (P1)"],
        [
            Step("read rows", "python_calamine.CalamineWorkbook.from_path; CalamineSheet.to_python", s_roundtrip),
            Step("seal", "json.dumps sort_keys; hashlib.sha256", lambda c: None),
            Step("compare", "csv_diff.compare", lambda c: None),
        ],
        "match, or the differing cell",
    )
    ctx = {"dir": out_dir("P30", label), "xlsx": OUT / "P16" / label / "facts.xlsx", "atoms": RESULTS[("facts", label)]["atoms"]}
    run(proof, label, ctx)
    return ctx


# ----------------------------------------------------------------------------
# P31 .. P40
# ----------------------------------------------------------------------------


def s_attendance_chart(ctx: dict[str, Any]) -> None:
    import plotly.graph_objects as go
    import plotly.io as pio

    rows = ctx["decision_rows"][1:]
    fig = go.Figure()
    fig.add_bar(x=[r[0] for r in rows], y=[r[1] for r in rows], name="voting members present")
    fig.add_hline(y=ctx["majority"], line_dash="dash", annotation_text=f"simple majority = {ctx['majority']}")
    fig.update_layout(title="TSC voting members present per meeting, against the charter's simple majority", xaxis_tickangle=-45)
    p = ctx["dir"] / "attendance.html"
    pio.write_html(fig, str(p), include_plotlyjs="cdn", full_html=True, div_id="chart")
    ctx["files"] = [p]
    ctx["shows"] = [f"{len(rows)} bars, threshold line at {ctx['majority']}"]


def attendance_chart() -> dict[str, Any]:
    proof = Proof(
        "P31",
        "decisions -> chart -> page",
        ["decisions (P13)"],
        [Step("chart", "plotly.graph_objects.Figure; add_bar; add_hline", s_attendance_chart), Step("page", "plotly.io.write_html(full_html, div_id)", lambda c: None)],
        "page",
    )
    src = RESULTS[("decisions", "majority")]
    ctx = {"dir": out_dir("P31", "majority"), "decision_rows": src["decision_rows"], "majority": src["majority"]}
    run(proof, "majority", ctx)
    return ctx


def s_coverage_matrix(ctx: dict[str, Any]) -> None:
    import duckdb

    con = duckdb.connect()
    con.execute("create table tag(meeting varchar, action varchar)")
    con.executemany("insert into tag values (?, ?)", [(t["sentence_id"].split(":")[0], ctx["lemma"].get(t["step"], "?")) for t in ctx["tags"] if not t["tie"]])
    counts = con.execute("select meeting, action, count(*) from tag group by meeting, action").fetchall()
    meetings = sorted({m for m, _, _ in counts})
    actions = sorted({a for _, a, _ in counts})
    cell = {(m, a): n for m, a, n in counts}
    rows = [["action"] + meetings] + [[a] + [cell.get((m, a), 0) for m in meetings] for a in actions]
    p = workbook(ctx["dir"] / "coverage_matrix.xlsx", {"matrix": rows})
    ctx["files"] = [p]
    ctx["shows"] = [f"matrix {len(actions)} actions x {len(meetings)} meetings", f"row {rows[1][0]}: {rows[1][1:]}"]


def coverage_matrix() -> dict[str, Any]:
    proof = Proof(
        "P32",
        "tagged actions -> coverage matrix -> workbook",
        ["tagged actions (P26)"],
        [Step("measure", "duckdb: count per (meeting, action); one row per action, one column per meeting", s_coverage_matrix), Step("tabulate", "openpyxl.Workbook.save", lambda c: None)],
        "coverage matrix",
    )
    src = RESULTS[("tagged steps", "nodejs-tsc-charter")]
    ctx = {"dir": out_dir("P32", "charter+minutes"), "tags": src["tags"], "lemma": src["lemma"]}
    run(proof, "charter+minutes", ctx)
    return ctx


def s_never_discussed(ctx: dict[str, Any]) -> None:
    import clingo

    prog = [f"step({lit(e)})." for e in ctx["order"]]
    prog += [f"tag({lit(t['sentence_id'])},{lit(t['step'])})." for t in ctx["tags"] if not t["tie"]]
    prog.append("discussed(E) :- tag(_,E).")
    prog.append("never(E) :- step(E), not discussed(E).")
    prog.append("lines(E,N) :- step(E), N = #count{ L : tag(L,E) }.")
    prog.append("#show never/1. #show lines/2.")
    ctl = clingo.Control(["0"])
    ctl.add("base", [], "\n".join(prog))
    ctl.ground([("base", [])])
    syms: list[Any] = []
    ctl.solve(on_model=lambda m: syms.extend(m.symbols(shown=True)))
    never = sorted(s.arguments[0].string for s in syms if s.name == "never")
    lines = {s.arguments[0].string: s.arguments[1].number for s in syms if s.name == "lines"}
    rows = [["step", "step_lemma", "minutes_lines", "never_discussed"]] + [[e, ctx["lemma"].get(e, "?"), lines.get(e, 0), e in never] for e in ctx["order"]]
    p = workbook(ctx["dir"] / "never_discussed.xlsx", {"steps": rows})
    ctx["files"] = [p]
    ctx["shows"] = [f"steps never discussed {len(never)} of {len(ctx['order'])}: {', '.join(ctx['lemma'].get(e, '?') for e in never)}"]


def never_discussed(label_steps: str) -> dict[str, Any]:
    proof = Proof(
        "P33",
        "ordered steps, tagged steps -> undiscussed steps -> workbook",
        ["ordered steps (P2), tagged steps (P10)"],
        [Step("complement", "clingo: never(E) :- step(E), not discussed(E)", s_never_discussed), Step("tabulate", "openpyxl.Workbook.save", lambda c: None)],
        "undiscussed steps",
    )
    src = RESULTS[("tagged steps", label_steps)]
    ctx = {"dir": out_dir("P33", label_steps), "tags": src["tags"], "lemma": src["lemma"], "order": src["order"]}
    run(proof, label_steps, ctx)
    return ctx


WOFLAN = r"""
import json, sys
import pm4py
from pm4py.objects.petri_net.importer import importer as pnml_importer
net, im, fm = pnml_importer.apply(sys.argv[1])
sound, diag = pm4py.check_soundness(net, im, fm)
keep = {str(getattr(k, "value", k)).split(".")[-1]: (v if isinstance(v, (bool, int, float, str)) else str(v)[:300]) for k, v in diag.items()}
print("RESULT " + json.dumps({"sound": bool(sound), "diagnostics": keep}))
"""
WOFLAN_BUDGET = 420


def s_soundness(ctx: dict[str, Any]) -> None:
    """woflan in a subprocess under a time budget; the verdict is memoized on the net's digest."""
    import subprocess

    import pm4py

    net, im, fm = ctx["net"]
    pnml = ctx["dir"] / "process.pnml"
    if not pnml.exists():
        pm4py.write_pnml(net, im, fm, str(pnml))
        canonical_xml(pnml)
    digest = sha256_file(pnml)
    out = ctx["dir"] / "soundness.json"
    res: dict[str, Any] | None = None
    if out.exists():
        prev = json.loads(out.read_text(encoding="utf-8"))
        if prev.get("pnml_sha256") == digest and prev.get("seconds_budget") == WOFLAN_BUDGET:
            res = prev
    if res is None:
        try:
            r = subprocess.run([sys.executable, "-c", WOFLAN, str(pnml)], capture_output=True, text=True, timeout=WOFLAN_BUDGET)
            line = [l for l in r.stdout.splitlines() if l.startswith("RESULT ")][-1]
            res = json.loads(line[len("RESULT "):])
        except subprocess.TimeoutExpired:
            res = {"sound": None, "undecided": f"woflan did not finish within {WOFLAN_BUDGET} s"}
        res["pnml_sha256"] = digest
        res["seconds_budget"] = WOFLAN_BUDGET
        write_json(out, res)
    ctx["files"] = [out] if pnml in ctx.get("files", []) else [pnml, out]
    verdict = res["sound"] if res.get("sound") is not None else res.get("undecided")
    scalars = [(k, v) for k, v in res.get("diagnostics", {}).items() if isinstance(v, (bool, int, float)) or (isinstance(v, str) and "\n" not in v and len(v) < 90)]
    ctx["shows"] = [f"sound: {verdict}", *[f"{k}: {v}" for k, v in scalars[:4]]]


def soundness_discovered(label: str) -> dict[str, Any]:
    proof = Proof(
        "P34",
        "process model -> soundness proof",
        ["process model (P8, discovered)"],
        [Step("soundness proof", "pm4py.check_soundness (woflan): workflow net, liveness, boundedness", s_soundness)],
        "soundness proof",
    )
    ctx = {"dir": out_dir("P34", label), "net": RESULTS[("log", label)]["net"]}
    run(proof, label, ctx)
    return ctx


def s_bpmn_to_net(ctx: dict[str, Any]) -> None:
    import pm4py

    bpmn = pm4py.read_bpmn(str(ctx["bpmn_path"]))
    net, im, fm = pm4py.convert_to_petri_net(bpmn)
    ctx["net"] = (net, im, fm)
    p = ctx["dir"] / "process.pnml"
    pm4py.write_pnml(net, im, fm, str(p))
    canonical_xml(p)
    ctx["files"] = [p]
    ctx["shows"] = [f"read back {len(bpmn.get_nodes())} bpmn nodes; petri net places {len(net.places)}, transitions {len(net.transitions)}"]


def s_soundness_append(ctx: dict[str, Any]) -> None:
    files, shows = ctx["files"], ctx["shows"]
    s_soundness(ctx)
    ctx["files"] = files + ctx["files"]
    ctx["shows"] = shows + ctx["shows"]


def soundness_from_order(label: str) -> dict[str, Any]:
    proof = Proof(
        "P35",
        "process -> petri net -> soundness proof",
        ["process (P5)"],
        [
            Step("read", "pm4py.read_bpmn", s_bpmn_to_net),
            Step("petri net", "pm4py.convert_to_petri_net; pm4py.write_pnml", lambda c: None),
            Step("soundness proof", "pm4py.check_soundness (woflan)", s_soundness_append),
        ],
        "soundness proof",
    )
    ctx = {"dir": out_dir("P35", label), "bpmn_path": RESULTS[("process", label)]["files"][0]}
    run(proof, label, ctx)
    return ctx


def s_variants(ctx: dict[str, Any]) -> None:
    import pm4py

    v = pm4py.get_variants_as_tuples(ctx["log"])
    items = sorted(((k, (n if isinstance(n, int) else len(n))) for k, n in v.items()), key=lambda kv: (-kv[1], kv[0]))
    rows = [["cases", "length", "variant"]] + [[n, len(k), " -> ".join(k)] for k, n in items]
    p = workbook(ctx["dir"] / "variants.xlsx", {"variants": rows})
    ctx["files"] = [p]
    ctx["shows"] = [f"variants {len(items)}; most common ({items[0][1]} cases): {' -> '.join(items[0][0])[:110]}"]


def variants(label: str) -> dict[str, Any]:
    proof = Proof(
        "P36",
        "log -> variants -> workbook",
        ["log (P8)"],
        [Step("variants", "pm4py.get_variants_as_tuples", s_variants), Step("tabulate", "openpyxl.Workbook.save", lambda c: None)],
        "variants",
    )
    ctx = {"dir": out_dir("P36", label), "log": RESULTS[("log", label)]["log"]}
    run(proof, label, ctx)
    return ctx


def s_rdf(ctx: dict[str, Any]) -> None:
    import rdflib
    from rdflib import Literal, Namespace, RDF, URIRef

    EX = Namespace("https://example.org/proofs/")
    g = rdflib.Graph()
    g.bind("ex", EX)

    def u(s: str) -> URIRef:
        return URIRef(EX + s.replace("#", "/"))

    for a in ctx["atoms"]:
        if a.predicate == "event":
            g.add((u(a.args[0]), RDF.type, EX.Event))
            g.add((u(a.args[0]), EX.lemma, Literal(a.args[1])))
            g.add((u(a.args[0]), EX.sentence, Literal(a.sentence_id)))
        elif a.predicate in ("agent", "patient", "theme"):
            g.add((u(a.args[0]), EX[a.predicate], u(a.args[1])))
            if (u(a.args[1]), EX.quote, None) not in g:
                g.add((u(a.args[1]), EX.quote, Literal(a.quote)))
        elif a.predicate in ("obligatory", "negated"):
            g.add((u(a.args[0]), EX[a.predicate], Literal(True)))
        elif a.predicate == "precedes":
            g.add((u(a.args[0]), EX.precedes, u(a.args[1])))
    p = ctx["dir"] / "facts.ttl"
    g.serialize(destination=str(p), format="turtle")
    q = """PREFIX ex: <https://example.org/proofs/>
SELECT ?e ?lemma ?who WHERE { ?e a ex:Event ; ex:lemma ?lemma ; ex:obligatory true ; ex:agent ?x . ?x ex:quote ?who } ORDER BY ?e"""
    res = [[str(r.e).replace(str(EX), ""), str(r.lemma), str(r.who)] for r in g.query(q)]
    ctx["sparql_count"] = len(res)
    rows = [["event", "lemma", "who"]] + res
    p2 = workbook(ctx["dir"] / "sparql_required_actions.xlsx", {"required": rows})
    (ctx["dir"] / "query.sparql").write_text(q + "\n", encoding="utf-8")
    ctx["files"] = [p, p2, ctx["dir"] / "query.sparql"]
    ctx["shows"] = [f"triples {len(g)}; obligatory events with an agent {len(res)}", *[f"{r[1]} by {r[2]}" for r in res[:3]]]


def rdf(label: str) -> dict[str, Any]:
    proof = Proof(
        "P37",
        "facts -> knowledge graph -> query results -> workbook",
        ["facts (P1)"],
        [
            Step("knowledge graph", "rdflib.Graph.add; Graph.serialize(format='turtle')", s_rdf),
            Step("query", "rdflib.Graph.query (SPARQL)", lambda c: None),
            Step("tabulate", "openpyxl.Workbook.save", lambda c: None),
        ],
        "query results",
    )
    ctx = {"dir": out_dir("P37", label), "atoms": RESULTS[("facts", label)]["atoms"]}
    run(proof, label, ctx)
    RESULTS[("rdf", label)] = ctx
    return ctx


def s_kuzu(ctx: dict[str, Any]) -> None:
    import shutil
    import tempfile

    import kuzu

    tmp = Path(tempfile.mkdtemp())
    db = kuzu.Database(str(tmp / "db"))
    con = kuzu.Connection(db)
    schema = [
        "CREATE NODE TABLE Event(id STRING, lemma STRING, obligatory BOOLEAN, negated BOOLEAN, sentence STRING, PRIMARY KEY(id))",
        "CREATE NODE TABLE Arg(id STRING, quote STRING, PRIMARY KEY(id))",
        "CREATE REL TABLE AGENT(FROM Event TO Arg)",
        "CREATE REL TABLE PATIENT(FROM Event TO Arg)",
        "CREATE REL TABLE THEME(FROM Event TO Arg)",
        "CREATE REL TABLE PRECEDES(FROM Event TO Event)",
    ]
    for s in schema:
        con.execute(s)
    atoms = ctx["atoms"]
    obl = {a.args[0] for a in atoms if a.predicate == "obligatory"}
    neg = {a.args[0] for a in atoms if a.predicate == "negated"}
    for a in atoms:
        if a.predicate == "event":
            con.execute("CREATE (:Event {id: $id, lemma: $lemma, obligatory: $o, negated: $n, sentence: $s})", {"id": a.args[0], "lemma": a.args[1], "o": a.args[0] in obl, "n": a.args[0] in neg, "s": a.sentence_id})
    seen: set[str] = set()
    for a in atoms:
        if a.predicate in ("agent", "patient", "theme") and a.args[1] not in seen:
            seen.add(a.args[1])
            con.execute("CREATE (:Arg {id: $id, quote: $q})", {"id": a.args[1], "q": a.quote})
    for a in atoms:
        if a.predicate in ("agent", "patient", "theme"):
            con.execute(f"MATCH (e:Event {{id: $e}}), (x:Arg {{id: $x}}) CREATE (e)-[:{a.predicate.upper()}]->(x)", {"e": a.args[0], "x": a.args[1]})
        elif a.predicate == "precedes":
            con.execute("MATCH (a:Event {id: $a}), (b:Event {id: $b}) CREATE (a)-[:PRECEDES]->(b)", {"a": a.args[0], "b": a.args[1]})
    q1 = "MATCH (e:Event)-[:AGENT]->(x:Arg) WHERE e.obligatory RETURN e.id, e.lemma, x.quote ORDER BY e.id"
    q2 = "MATCH (a:Event)-[:PRECEDES]->(b:Event) RETURN a.lemma, b.lemma, a.id, b.id ORDER BY a.id, b.id"
    r1 = con.execute(q1).get_all()
    r2 = con.execute(q2).get_all()
    p = workbook(ctx["dir"] / "cypher_results.xlsx", {"required_with_agent": [["event", "lemma", "who"]] + r1, "precedes": [["before", "after", "before_id", "after_id"]] + r2})
    (ctx["dir"] / "queries.cypher").write_text("\n".join(schema + [q1, q2]) + "\n", encoding="utf-8")
    shutil.rmtree(tmp, ignore_errors=True)
    ctx["files"] = [p, ctx["dir"] / "queries.cypher"]
    sparql = ctx.get("sparql_count")
    ctx["shows"] = [f"events {len([a for a in atoms if a.predicate == 'event'])}, args {len(seen)}; obligatory events with an agent {len(r1)}; precedes edges {len(r2)}", f"same count as the SPARQL query of P37: {sparql == len(r1)} ({sparql})", *[f"{r[1]} by {r[2]}" for r in r1[:2]], *[f"{r[0]} precedes {r[1]}" for r in r2[:2]]]


def graph_db(label: str) -> dict[str, Any]:
    proof = Proof(
        "P38",
        "facts -> graph database -> query results -> workbook",
        ["facts (P1)"],
        [
            Step("graph database", "kuzu.Database; kuzu.Connection.execute: CREATE NODE TABLE, CREATE REL TABLE, CREATE", s_kuzu),
            Step("query", "kuzu.Connection.execute (Cypher MATCH); QueryResult.get_all", lambda c: None),
            Step("tabulate", "openpyxl.Workbook.save", lambda c: None),
        ],
        "query results",
    )
    ctx = {"dir": out_dir("P38", label), "atoms": RESULTS[("facts", label)]["atoms"], "sparql_count": RESULTS.get(("rdf", label), {}).get("sparql_count")}
    run(proof, label, ctx)
    return ctx


def s_site(ctx: dict[str, Any]) -> None:
    import shutil

    from mkdocs.commands.build import build
    from mkdocs.config import load_config

    d = ctx["dir"]
    docs = d / "docs"
    if docs.exists():
        shutil.rmtree(docs)
    docs.mkdir(parents=True)
    o = ctx["ordering"]
    atoms = ctx["atoms"]
    lemma = {a.args[0]: a.args[1] for a in atoms if a.predicate == "event"}
    sent = {ps.sentence.id: ps.sentence.text for ps in ctx["parsed"]}
    idx = [f"# {ctx['label']}: ordered steps", "", "| # | step | sentence |", "|---|---|---|"]
    bar = "\\|"
    for i, e in enumerate(o.order, 1):
        text = sent.get(e.split("#")[0], "").replace("|", bar)
        idx.append(f"| {i} | {lemma.get(e, '?')} `{e.split(':', 1)[1]}` | {text} |")
    idx += ["", "## forced precedence", "", "| before | after |", "|---|---|"] + [f"| {lemma.get(a, '?')} | {lemma.get(b, '?')} |" for a, b in o.forced]
    (docs / "index.md").write_text("\n".join(idx) + "\n", encoding="utf-8")
    f = ["# facts", "", "| predicate | args | quote |", "|---|---|---|"]
    for a in atoms:
        args = " ".join(x.split("#")[-1] for x in a.args)
        f.append(f"| {a.predicate} | {args} | {a.quote.replace('|', bar)} |")
    (docs / "facts.md").write_text("\n".join(f) + "\n", encoding="utf-8")
    cfg = d / "mkdocs.yml"
    cfg.write_text(f"site_name: {ctx['label']}\ndocs_dir: docs\nsite_dir: site\nuse_directory_urls: false\nnav:\n  - steps: index.md\n  - facts: facts.md\n", encoding="utf-8")
    build(load_config(str(cfg)))
    site = d / "site"
    for junk in ("sitemap.xml", "sitemap.xml.gz"):
        (site / junk).unlink(missing_ok=True)
    ctx["files"] = [site / "index.html", site / "facts.html"]
    ctx["shows"] = [f"pages {len(list(site.glob('*.html')))}; index rows {len(o.order)}; facts rows {len(atoms)}"]


def site(label: str) -> dict[str, Any]:
    proof = Proof(
        "P39",
        "ordered steps, facts -> pages -> site",
        ["ordered steps (P2), facts (P1)"],
        [
            Step("pages", "markdown tables written from the rows", s_site),
            Step("site", "mkdocs.config.load_config; mkdocs.commands.build.build", lambda c: None),
        ],
        "site",
    )
    src = RESULTS[("ordered steps", label)]
    ctx = {"dir": out_dir("P39", label), "ordering": src["ordering"], "atoms": src["atoms"], "parsed": src["parsed"], "label": label}
    run(proof, label, ctx)
    return ctx


def s_doc_roundtrip(ctx: dict[str, Any]) -> None:
    import docx
    from csv_diff import compare

    d = docx.Document(str(ctx["docx"]))
    back = []
    for para in d.paragraphs:
        m = re.match(r"^(\d+)\. (.+?)  \[(.+)\]$", para.text)
        if m:
            back.append({"position": m.group(1), "lemma": m.group(2), "id": m.group(3)})
    o = ctx["ordering"]
    lemma = {a.args[0]: a.args[1] for a in ctx["atoms"] if a.predicate == "event"}
    fwd = [{"position": str(i), "lemma": lemma.get(e, "?"), "id": e} for i, e in enumerate(o.order, 1)]
    diff = compare({r["id"]: r for r in fwd}, {r["id"]: r for r in back})
    same = sha256_json(fwd) == sha256_json(back)
    p = write_json(ctx["dir"] / "roundtrip.json", {"digest_forward": sha256_json(fwd), "digest_back": sha256_json(back), "match": same, "changed": diff["changed"][:5], "added": diff["added"][:5], "removed": diff["removed"][:5]})
    ctx["files"] = [p]
    RESULTS.setdefault(("roundtrips", "all"), []).append({"proof": "P40", "label": "document", "match": bool(same), "changed": len(diff["changed"])})
    ctx["shows"] = [f"steps read back {len(back)} of {len(fwd)}; match {same}; changed {len(diff['changed'])}, added {len(diff['added'])}, removed {len(diff['removed'])}"]


def document_roundtrip(label: str) -> dict[str, Any]:
    proof = Proof(
        "P40",
        "document -> steps' -> digest' ; digest', digest -> match",
        ["document (P15), ordered steps (P2)"],
        [
            Step("read paragraphs", "docx.Document; Document.paragraphs", s_doc_roundtrip),
            Step("seal", "json.dumps sort_keys; hashlib.sha256", lambda c: None),
            Step("compare", "csv_diff.compare", lambda c: None),
        ],
        "match, or the differing cell",
    )
    src = RESULTS[("ordered steps", label)]
    ctx = {"dir": out_dir("P40", label), "docx": OUT / "P15" / label / "ordered_steps.docx", "ordering": src["ordering"], "atoms": src["atoms"]}
    run(proof, label, ctx)
    return ctx


# ----------------------------------------------------------------------------
# P41 .. P50
# ----------------------------------------------------------------------------


def s_bpmn_3d(ctx: dict[str, Any]) -> None:
    import networkx as nx
    import plotly.graph_objects as go
    import plotly.io as pio

    o = ctx["ordering"]
    lemma = {a.args[0]: a.args[1] for a in ctx["atoms"] if a.predicate == "event"}
    G = nx.DiGraph()
    G.add_nodes_from(o.order)
    G.add_edges_from(o.forced)
    layers = list(nx.topological_generations(G))
    pos: dict[str, tuple[int, int, int]] = {}
    for li, layer in enumerate(layers):
        for ni, n in enumerate(sorted(layer)):
            unit = int(re.search(r":u(\d+):", n).group(1))
            pos[n] = (li, ni, unit)
    fig = go.Figure()
    for a, b in o.forced:
        fig.add_trace(go.Scatter3d(x=[pos[a][0], pos[b][0]], y=[pos[a][1], pos[b][1]], z=[pos[a][2], pos[b][2]], mode="lines", line={"color": "gray"}, showlegend=False))
    fig.add_trace(
        go.Scatter3d(
            x=[pos[n][0] for n in o.order], y=[pos[n][1] for n in o.order], z=[pos[n][2] for n in o.order], mode="markers+text",
            text=[lemma.get(n, "?") for n in o.order], textposition="top center", marker={"size": 6}, showlegend=False,
        )
    )
    fig.update_layout(title=f"{ctx['label']}: forced order in three axes", scene={"xaxis_title": "layer (earlier to later)", "yaxis_title": "position in layer", "zaxis_title": "unit of the source text"})
    p = ctx["dir"] / "process_3d.html"
    pio.write_html(fig, str(p), include_plotlyjs="cdn", full_html=True, div_id="chart")
    ctx["files"] = [p]
    ctx["shows"] = [f"nodes {len(pos)}, edges {len(o.forced)}, layers {len(layers)}; z = source unit index"]


def bpmn_3d(label: str) -> dict[str, Any]:
    proof = Proof(
        "P41",
        "ordered steps -> layers -> 3D process page",
        ["ordered steps (P2)"],
        [
            Step("layers", "networkx.topological_generations; z from the source unit of each step", s_bpmn_3d),
            Step("draw", "plotly.graph_objects.Scatter3d (lines, markers+text)", lambda c: None),
            Step("page", "plotly.io.write_html(full_html, div_id)", lambda c: None),
        ],
        "3D process page",
    )
    src = RESULTS[("ordered steps", label)]
    ctx = {"dir": out_dir("P41", label), "ordering": src["ordering"], "atoms": src["atoms"], "label": label}
    run(proof, label, ctx)
    return ctx


POLICY_PAGE = """<!doctype html><html><head><meta charset="utf-8"><title>{{ title }}</title>
<style>body{font-family:system-ui,sans-serif;max-width:60rem;margin:2rem auto;padding:0 1rem;color:#222}
table{border-collapse:collapse;width:100%}td,th{border:1px solid #ccc;padding:.4rem .6rem;vertical-align:top;text-align:left}
th{background:#f4f4f4}.rule{color:#555;font-size:.9em}.yes{color:#0a6}.no{color:#c33}</style></head><body>
<h1>{{ title }}</h1>
<p>{{ subtitle }}</p>
<table><tr><th>if</th><th>then</th><th>because</th></tr>
{% for r in rules %}<tr><td>{{ r.if }}</td><td>{{ r.then }}</td><td class="rule">{{ r.because }}</td></tr>
{% endfor %}</table>
{% if cases %}<h2>applied</h2><table><tr>{% for h in case_header %}<th>{{ h }}</th>{% endfor %}</tr>
{% for c in cases %}<tr>{% for v in c %}<td class="{{ 'yes' if v is true else ('no' if v is false else '') }}">{{ v }}</td>{% endfor %}</tr>
{% endfor %}</table>{% endif %}
</body></html>
"""


def s_policy_page(ctx: dict[str, Any]) -> None:
    import jinja2

    jdm = json.loads(Path(ctx["jdm_path"]).read_text(encoding="utf-8"))
    dt = next(n for n in jdm["nodes"] if n["type"] == "decisionTableNode")
    inputs = {i["id"]: i["name"] for i in dt["content"]["inputs"]}
    outputs = {o["id"]: o["name"] for o in dt["content"]["outputs"]}
    rules = []
    for r in dt["content"]["rules"]:
        conds = [f"{inputs[k]} {v}" for k, v in r.items() if k in inputs and v]
        outs = [f"{outputs[k]} = {v.strip(chr(39))}" for k, v in r.items() if k in outputs]
        rules.append({"if": " and ".join(conds) or "otherwise", "then": "; ".join(outs), "because": r.get("_description", "")})
    html = jinja2.Environment(autoescape=True).from_string(POLICY_PAGE).render(
        title=ctx["title"], subtitle=ctx["subtitle"], rules=rules, cases=ctx.get("cases", []), case_header=ctx.get("case_header", [])
    )
    p = ctx["dir"] / "policy.html"
    p.write_text(html, encoding="utf-8")
    ctx["files"] = [p]
    ctx["shows"] = [f"rules {len(rules)}; first: if {rules[0]['if']} then {rules[0]['then']}"]


def policy_page(label: str) -> dict[str, Any]:
    proof = Proof(
        "P42",
        "policy -> readable policy page",
        ["policy (P4 or P13), decisions (P13)"],
        [
            Step("read table", "json.loads of the GoRules JDM; inputs, outputs, rules, rule descriptions", s_policy_page),
            Step("page", "jinja2.Environment(autoescape).from_string(...).render -> html", lambda c: None),
        ],
        "policy page",
    )
    if label == "majority":
        src = RESULTS[("decisions", "majority")]
        ctx = {
            "dir": out_dir("P42", label), "jdm_path": src["files"][0], "title": "Can a vote carry at this meeting?",
            "subtitle": "The rule is the charter's sentence, rendered as a table; the cases are the eight meetings' attendance.",
            "cases": src["decision_rows"][1:], "case_header": src["decision_rows"][0],
        }
    else:
        src = RESULTS[("policy", label)]
        ctx = {"dir": out_dir("P42", label), "jdm_path": src["files"][0], "title": f"{label}: required actions", "subtitle": "Each row is an obligation found in the text; the quote is the sentence it came from."}
    run(proof, label, ctx)
    return ctx


def s_heatmap(ctx: dict[str, Any]) -> None:
    import plotly.express as px
    import plotly.io as pio

    con = duck_events(ctx["rows"])
    res = con.execute(
        """
        with seq as (select case_id, activity, department, ts, lag(ts) over (partition by case_id order by ts) as prev from ev)
        select department, activity, round(avg(epoch(ts - prev))/3600, 2) as mean_hours
        from seq where prev is not null group by department, activity order by department, activity
        """
    ).fetchall()
    deps = sorted({r[0] for r in res})
    acts = sorted({r[1] for r in res})
    cell = {(d, a): h for d, a, h in res}
    z = [[cell.get((d, a)) for a in acts] for d in deps]
    fig = px.imshow(z, x=acts, y=deps, labels={"x": "activity", "y": "department", "color": "mean hours since previous event"}, title="waiting time by department and activity, receipt phase", aspect="auto")
    fig.update_layout(xaxis_tickangle=-45)
    p = ctx["dir"] / "heatmap.html"
    pio.write_html(fig, str(p), include_plotlyjs="cdn", full_html=True, div_id="chart")
    p2 = workbook(ctx["dir"] / "heatmap.xlsx", {"mean_hours": [["department"] + acts] + [[d] + row for d, row in zip(deps, z)]})
    ctx["files"] = [p, p2]
    ctx["shows"] = [f"{len(deps)} departments x {len(acts)} activities", *[f"{r[0]} / {r[1]}: {r[2]} h" for r in sorted(res, key=lambda r: -(r[2] or 0))[:3]]]


def heatmap(label: str) -> dict[str, Any]:
    proof = Proof(
        "P43",
        "records -> measured cells -> heatmap page",
        ["records (P6)"],
        [
            Step("measure", "duckdb: avg(epoch(ts - lag(ts))) per department and activity", s_heatmap),
            Step("chart", "plotly.express.imshow", lambda c: None),
            Step("page", "plotly.io.write_html(full_html, div_id)", lambda c: None),
        ],
        "heatmap page",
    )
    ctx = {"dir": out_dir("P43", label), "rows": RESULTS[("records", label)]["rows"]}
    run(proof, label, ctx)
    return ctx


def s_who_what(ctx: dict[str, Any]) -> None:
    import duckdb

    atoms = ctx["atoms"]
    lemma = {a.args[0]: a.args[1] for a in atoms if a.predicate == "event"}
    obl = {a.args[0] for a in atoms if a.predicate == "obligatory"}
    con = duckdb.connect()
    con.execute("create table t(who varchar, what varchar)")
    con.executemany("insert into t values (?, ?)", [(a.quote.lower(), lemma.get(a.args[0], "?")) for a in atoms if a.predicate == "agent" and a.args[0] in obl])
    res = con.execute("select who, what, count(*) as n from t group by who, what order by n desc, who, what").fetchall()
    whos = sorted({r[0] for r in res})
    whats = sorted({r[1] for r in res})
    cell = {(w, x): n for w, x, n in res}
    rows = [["who \\ must"] + whats] + [[w] + [cell.get((w, x), 0) for x in whats] for w in whos]
    p = workbook(ctx["dir"] / "who_must_what.xlsx", {"matrix": rows, "pairs": [["who", "must", "count"]] + [list(r) for r in res]})
    ctx["files"] = [p]
    ctx["shows"] = [f"{len(whos)} actors x {len(whats)} required actions", *[f"{r[0]} must {r[1]} ({r[2]})" for r in res[:3]]]


def who_must_what(label: str) -> dict[str, Any]:
    proof = Proof(
        "P44",
        "facts -> actor by action matrix -> workbook",
        ["facts (P1)"],
        [Step("measure", "duckdb: count per (agent quote, event lemma) over obligatory events", s_who_what), Step("tabulate", "openpyxl.Workbook.save", lambda c: None)],
        "actor by action matrix",
    )
    ctx = {"dir": out_dir("P44", label), "atoms": RESULTS[("facts", label)]["atoms"]}
    run(proof, label, ctx)
    return ctx


def s_timeline(ctx: dict[str, Any]) -> None:
    import plotly.graph_objects as go
    import plotly.io as pio

    pts = collections.Counter()
    for t in ctx["tags"]:
        if not t["tie"]:
            pts[(t["sentence_id"].split(":")[0].replace("tsc-", ""), ctx["lemma"].get(t["step"], "?"))] += 1
    xs = [k[0] for k in pts]
    ys = [k[1] for k in pts]
    sizes = [6 + 3 * n for n in pts.values()]
    fig = go.Figure(go.Scatter(x=xs, y=ys, mode="markers", marker={"size": sizes}, text=[f"{n} lines" for n in pts.values()]))
    fig.update_layout(title="which governance steps each TSC meeting discussed", xaxis_title="meeting", yaxis_title="step")
    p = ctx["dir"] / "timeline.html"
    pio.write_html(fig, str(p), include_plotlyjs="cdn", full_html=True, div_id="chart")
    ctx["files"] = [p]
    ctx["shows"] = [f"points {len(pts)}; meetings {len(set(xs))}, steps {len(set(ys))}"]


def timeline(label_steps: str) -> dict[str, Any]:
    proof = Proof(
        "P45",
        "tagged steps -> timeline page",
        ["tagged steps (P10)"],
        [Step("measure", "count of lines per (meeting, step)", s_timeline), Step("chart", "plotly.graph_objects.Scatter (marker size = lines)", lambda c: None), Step("page", "plotly.io.write_html", lambda c: None)],
        "timeline page",
    )
    src = RESULTS[("tagged steps", label_steps)]
    ctx = {"dir": out_dir("P45", label_steps), "tags": src["tags"], "lemma": src["lemma"]}
    run(proof, label_steps, ctx)
    return ctx


def s_bottlenecks(ctx: dict[str, Any]) -> None:
    import duckdb

    con = duckdb.connect()
    con.execute("create table dfg(a varchar, b varchar, n integer)")
    con.executemany("insert into dfg values (?, ?, ?)", [tuple(r) for r in ctx["dfg_rows"][1:]])
    con.execute("create table wait(activity varchar, events integer, mean_hours double)")
    con.executemany("insert into wait values (?, ?, ?)", [(r[0], r[1], r[2]) for r in ctx["measured"][1:]])
    res = con.execute(
        "select d.a, d.b, d.n, w.mean_hours from dfg d join wait w on w.activity = d.b where w.mean_hours is not null order by w.mean_hours desc, d.n desc limit 25"
    ).fetchall()
    rows = [["from", "to", "cases", "mean_hours_waiting_into_to"]] + [list(r) for r in res]
    RESULTS[("bottlenecks", "all")] = [dict(zip(rows[0], r)) for r in rows[1:]]
    p = workbook(ctx["dir"] / "bottlenecks.xlsx", {"bottlenecks": rows})
    ctx["files"] = [p]
    ctx["shows"] = [f"{r[0]} -> {r[1]}: {r[2]} cases, {r[3]} h" for r in res[:3]]


def bottlenecks(label: str) -> dict[str, Any]:
    proof = Proof(
        "P46",
        "directly-follows graph, measured steps -> bottlenecks -> workbook",
        ["directly-follows graph (P23), measured steps (P12)"],
        [Step("key", "duckdb: join edges to the waiting time of their target activity", s_bottlenecks), Step("rank", "duckdb: order by mean hours desc", lambda c: None), Step("tabulate", "openpyxl.Workbook.save", lambda c: None)],
        "bottlenecks",
    )
    ctx = {"dir": out_dir("P46", label), "dfg_rows": RESULTS[("dfg", label)]["dfg_rows"], "measured": RESULTS[("measured activities", "receipt-xes+receipt-csv")]["measured"]}
    run(proof, label, ctx)
    return ctx


AT_RISK_JDM = """{
 "contentType": "application/vnd.gorules.decision",
 "nodes": [
  {"id": "in", "type": "inputNode", "name": "case", "position": {"x": 0, "y": 0}},
  {"id": "dt", "type": "decisionTableNode", "name": "open case past its department's mean", "position": {"x": 0, "y": 0},
   "content": {"hitPolicy": "first",
    "inputs": [{"id": "i1", "name": "department", "field": "department"}, {"id": "i2", "name": "elapsed days", "field": "elapsed_days"}],
    "outputs": [{"id": "o1", "name": "at risk", "field": "at_risk"}, {"id": "o2", "name": "rule", "field": "rule"}],
    "rules": [
{% for d in deps %}     {"_id": "r{{ loop.index }}", "_description": "open longer than the {{ d.department }} mean of {{ d.mean_days }} days (P22)", "i1": {{ (d.department|tojson)|tojson }}, "i2": "> {{ d.mean_days }}", "o1": "true", "o2": "'r{{ loop.index }}'"},
{% endfor %}     {"_id": "r0", "_description": "otherwise", "i1": "", "i2": "", "o1": "false", "o2": "'r0'"}
    ]}},
  {"id": "out", "type": "outputNode", "name": "decision", "position": {"x": 0, "y": 0}}
 ],
 "edges": [{"id": "e1", "sourceId": "in", "targetId": "dt", "type": "edge"}, {"id": "e2", "sourceId": "dt", "targetId": "out", "type": "edge"}]
}
"""


def s_at_risk(ctx: dict[str, Any]) -> None:
    import datetime as _dt

    import jinja2
    import zen

    con = duck_events(ctx["rows"])
    deps = con.execute(
        "with c as (select case_id, any_value(department) d, min(ts) f, max(ts) l from ev group by case_id) select d, round(avg(epoch(l - f))/86400, 2) from c group by d order by d"
    ).fetchall()
    end = con.execute("select max(ts) from ev").fetchone()[0]
    opens = con.execute(
        "with c as (select case_id, any_value(department) d, any_value(enddate) e, max(ts) l from ev group by case_id) select case_id, d, round(epoch(? - l)/86400, 2) from c where e = '' order by case_id", [end]
    ).fetchall()
    jdm = jinja2.Environment().from_string(AT_RISK_JDM).render(deps=[{"department": d, "mean_days": m} for d, m in deps])
    p = ctx["dir"] / "policy.jdm.json"
    p.write_text(jdm, encoding="utf-8")
    dec = zen.ZenEngine().create_decision(jdm)
    rows = [["case", "department", "elapsed_days_at_log_end", "at_risk", "rule"]]
    for c, d, el in opens:
        r = dec.evaluate({"department": d, "elapsed_days": el})["result"]
        rows.append([c, d, el, r["at_risk"], r["rule"]])
    p2 = workbook(ctx["dir"] / "at_risk.xlsx", {"open_cases": rows})
    n = sum(1 for r in rows[1:] if r[3])
    ctx["files"] = [p, p2]
    ctx["shows"] = [f"log ends {end.date() if hasattr(end, 'date') else end}; open cases {len(opens)}; at risk {n}", f"department means {deps[:3]}", f"first rows {rows[1:3]}"]


def at_risk(label: str) -> dict[str, Any]:
    proof = Proof(
        "P47",
        "records, measured groups -> decision table -> decisions -> workbook",
        ["records (P6), measured groups (P22)"],
        [
            Step("measure", "duckdb: department mean duration; elapsed days of each open case at the log's end", s_at_risk),
            Step("decision table", "jinja2 render: one rule per department, threshold = its mean -> GoRules JDM", lambda c: None),
            Step("evaluate", "zen.ZenEngine.create_decision; zen.ZenDecision.evaluate per open case", lambda c: None),
            Step("tabulate", "openpyxl.Workbook.save", lambda c: None),
        ],
        "decisions",
    )
    ctx = {"dir": out_dir("P47", label), "rows": RESULTS[("records", label)]["rows"]}
    run(proof, label, ctx)
    return ctx


def s_checklist(ctx: dict[str, Any]) -> None:
    import docx

    sent = {ps.sentence.id: ps.sentence.text for ps in ctx["parsed"]}
    d = docx.Document()
    d.add_heading(f"{ctx['label']}: checklist of required actions", 1)
    seen: set[str] = set()
    n = 0
    for r in ctx["required"]:
        key = (r["sentence_id"], r["action"])
        if key in seen:
            continue
        seen.add(key)
        d.add_paragraph(f"☐  {r['action']}")
        d.add_paragraph(sent.get(r["sentence_id"], ""), style="Intense Quote")
        n += 1
    p = ctx["dir"] / "checklist.docx"
    d.save(p)
    ctx["files"] = [p]
    ctx["shows"] = [f"checklist items {n}"]


def checklist(label: str) -> dict[str, Any]:
    proof = Proof(
        "P48",
        "required actions -> checklist document",
        ["required actions (P4), facts (P1)"],
        [Step("write", "docx.Document; add_paragraph with a box glyph per required action and its sentence; docx.document.Document.save", s_checklist)],
        "checklist document",
    )
    ctx = dict(RESULTS[("policy", label)])
    ctx["dir"] = out_dir("P48", label)
    run(proof, label, ctx)
    return ctx


def s_cross_document(ctx: dict[str, Any]) -> None:
    import clingo

    prog = [f"req({lit(lbl)},{lit(r['action'])})." for lbl, reqs in ctx["required"].items() for r in reqs]
    prog.append("both(A) :- req(D1,A), req(D2,A), D1 < D2.")
    prog.append("only(D,A) :- req(D,A), not both(A).")
    prog.append("#show both/1. #show only/2.")
    ctl = clingo.Control(["0"])
    ctl.add("base", [], "\n".join(prog))
    ctl.ground([("base", [])])
    syms: list[Any] = []
    ctl.solve(on_model=lambda m: syms.extend(m.symbols(shown=True)))
    both = sorted(s.arguments[0].string for s in syms if s.name == "both")
    only = sorted((s.arguments[0].string, s.arguments[1].string) for s in syms if s.name == "only")
    rows = [["action", "in"]] + [[a, "both"] for a in both] + [[a, d] for d, a in only]
    p = workbook(ctx["dir"] / "cross_document.xlsx", {"required_actions": rows})
    ctx["files"] = [p]
    ctx["shows"] = [f"required in both documents: {both}", f"only in one: {len(only)}"]


def cross_document(labels: list[str]) -> dict[str, Any]:
    proof = Proof(
        "P49",
        "required actions, required actions -> shared and unshared actions -> workbook",
        ["required actions of two documents (P4)"],
        [Step("key", "clingo: join on the action lemma across documents", s_cross_document), Step("tabulate", "openpyxl.Workbook.save", lambda c: None)],
        "shared and unshared actions",
    )
    ctx = {"dir": out_dir("P49", "+".join(labels)), "required": {l: RESULTS[("policy", l)]["required"] for l in labels}}
    run(proof, "+".join(labels), ctx)
    return ctx


def s_site_roundtrip(ctx: dict[str, Any]) -> None:
    import lxml.html
    from csv_diff import compare

    doc = lxml.html.fromstring(Path(ctx["index"]).read_text(encoding="utf-8"))
    table = doc.xpath("//table")[0]
    back = []
    for tr in table.xpath(".//tr")[1:]:
        cells = [" ".join(td.text_content().split()) for td in tr.xpath("./td")]
        m = re.match(r"^(.+?) (\S+)$", cells[1])
        back.append({"position": cells[0], "lemma": m.group(1), "id": m.group(2)})
    o = ctx["ordering"]
    lemma = {a.args[0]: a.args[1] for a in ctx["atoms"] if a.predicate == "event"}
    fwd = [{"position": str(i), "lemma": lemma.get(e, "?"), "id": e.split(":", 1)[1]} for i, e in enumerate(o.order, 1)]
    diff = compare({r["id"]: r for r in fwd}, {r["id"]: r for r in back})
    same = sha256_json(fwd) == sha256_json(back)
    p = write_json(ctx["dir"] / "roundtrip.json", {"digest_forward": sha256_json(fwd), "digest_back": sha256_json(back), "match": same, "changed": diff["changed"][:5], "added": diff["added"][:5], "removed": diff["removed"][:5]})
    ctx["files"] = [p]
    RESULTS.setdefault(("roundtrips", "all"), []).append({"proof": "P50", "label": "site", "match": bool(same), "changed": len(diff["changed"])})
    ctx["shows"] = [f"rows read back {len(back)} of {len(fwd)}; match {same}; changed {len(diff['changed'])}"]


def site_roundtrip(label: str) -> dict[str, Any]:
    proof = Proof(
        "P50",
        "site -> steps' -> digest' ; digest', digest -> match",
        ["site (P39), ordered steps (P2)"],
        [
            Step("read page", "lxml.html.fromstring; xpath over the steps table", s_site_roundtrip),
            Step("seal", "json.dumps sort_keys; hashlib.sha256", lambda c: None),
            Step("compare", "csv_diff.compare", lambda c: None),
        ],
        "match, or the differing cell",
    )
    src = RESULTS[("ordered steps", label)]
    ctx = {"dir": out_dir("P50", label), "index": OUT / "P39" / label / "site" / "index.html", "ordering": src["ordering"], "atoms": src["atoms"]}
    run(proof, label, ctx)
    return ctx


# ----------------------------------------------------------------------------
# the ITIL baseline (proofs/itil.yaml through proofs/itil.py)
# ----------------------------------------------------------------------------

ITIL_REPORT: dict[str, Any] = {}


def bpmn_from_edges(rows: list[dict[str, Any]], path: Path) -> Path:
    """A BPMN with parallel gateways from rows (before, after, before_step, after_step)."""
    from pm4py.objects.bpmn.exporter import exporter as bpmn_exporter
    from pm4py.objects.bpmn.obj import BPMN

    edges = [(r["before"], r["after"]) for r in rows]
    names = {}
    for r in rows:
        names[r["before"]] = r.get("before_step", "?")
        names[r["after"]] = r.get("after_step", "?")
    order = sorted(names)
    preds, succs = collections.defaultdict(set), collections.defaultdict(set)
    for a, c in edges:
        preds[c].add(a)
        succs[a].add(c)
    b = BPMN()
    start, end = BPMN.StartEvent(name="start"), BPMN.EndEvent(name="end")
    b.add_node(start)
    b.add_node(end)
    entry, exit_ = {}, {}
    for e in order:
        t = BPMN.Task(name=f"{names[e]} [{e.split(':', 1)[1] if ':' in e else e}]")
        b.add_node(t)
        entry[e], exit_[e] = t, t
        if len(preds[e]) > 1:
            g = BPMN.ParallelGateway(name=f"join {e}", gateway_direction=BPMN.Gateway.Direction.CONVERGING)
            b.add_node(g)
            b.add_flow(BPMN.SequenceFlow(g, t))
            entry[e] = g
        if len(succs[e]) > 1:
            g = BPMN.ParallelGateway(name=f"split {e}", gateway_direction=BPMN.Gateway.Direction.DIVERGING)
            b.add_node(g)
            b.add_flow(BPMN.SequenceFlow(t, g))
            exit_[e] = g
    for a, c in edges:
        b.add_flow(BPMN.SequenceFlow(exit_[a], entry[c]))
    sources = [e for e in order if not preds[e]]
    sinks = [e for e in order if not succs[e]]
    first, last = start, end
    if len(sources) > 1:
        first = BPMN.ParallelGateway(name="split start", gateway_direction=BPMN.Gateway.Direction.DIVERGING)
        b.add_node(first)
        b.add_flow(BPMN.SequenceFlow(start, first))
    if len(sinks) > 1:
        last = BPMN.ParallelGateway(name="join end", gateway_direction=BPMN.Gateway.Direction.CONVERGING)
        b.add_node(last)
        b.add_flow(BPMN.SequenceFlow(last, end))
    for e in sources:
        b.add_flow(BPMN.SequenceFlow(first, entry[e]))
    for e in sinks:
        b.add_flow(BPMN.SequenceFlow(exit_[e], last))
    bpmn_exporter.apply(b, str(path))
    canonical_xml(path)
    return path


def run_itil() -> dict[str, Any]:
    sys.path.insert(0, str(ROOT / "proofs"))
    import itil

    helpers = {
        "traces": RESULTS.get(("traces", "all"), {}),
        "roundtrips": RESULTS.get(("roundtrips", "all"), []),
        "bottlenecks": RESULTS.get(("bottlenecks", "all"), []),
        "bpmn_from_edges": bpmn_from_edges,
    }
    ctx = itil.Ctx(RESULTS, helpers)
    report = itil.run_register(ctx, ROOT / "proofs" / "itil.yaml")
    write_json(OUT / "itil" / "report.json", report)
    return report


def main() -> None:
    for label in ["usc5-552-doj", "nodejs-tsc-charter", "nodejs-governance"]:
        door_one(label)
    for label in ["usc5-552-doj", "nodejs-tsc-charter", "nodejs-governance"]:
        door_one_order(label)
    anchor_dates("usc5-552-doj")
    anchor_dates("nodejs-tsc-charter")
    policy("usc5-552-doj")
    policy("nodejs-tsc-charter")
    process_from_order("usc5-552-doj")
    process_from_order("nodejs-governance")
    door_two("receipt-csv", ROOT / "proofs" / "in" / "receipt.csv", s_read_rows_csv, "receipt_event", "event_row", ["case:concept:name", "concept:name", "time:timestamp", "org:resource", "case:department", "case:channel"])
    door_two("tsc-voting-members", ROOT / "proofs" / "in" / "node-README.md", s_read_rows_roster, "voting_member", "voting_member", ["handle", "name"])
    door_three("nodejs-tsc-minutes")
    door_four("receipt-xes", ROOT / "proofs" / "in" / "receipt.xes")
    replay("receipt-xes")
    tag("nodejs-governance", "nodejs-tsc-minutes")
    measure_tags("nodejs-governance")
    key_measure_log("receipt-xes", "receipt-csv")
    evaluate_majority()
    deck("usc5-552-doj")
    deck("nodejs-governance")
    document("usc5-552-doj")
    tabulate("usc5-552-doj")
    compose("nodejs-governance", "tsc-2024-01-17")
    reverse("nodejs-governance")
    minutes_as_log("nodejs-governance")
    document("nodejs-governance")
    for label in ["nodejs-tsc-charter", "nodejs-governance"]:
        tabulate(label)
    deadlines("receipt-csv")
    by_department("receipt-csv")
    dfg("receipt-xes")
    dashboard("receipt-xes+receipt-csv")
    obligations_document("usc5-552-doj")
    obligations_document("nodejs-tsc-charter")
    tag("nodejs-tsc-charter", "nodejs-tsc-minutes", pid="P26", required=True)
    measure_tags("nodejs-tsc-charter")
    conformance_by_department("receipt-xes", "receipt-csv")
    decisions_with_discussion()
    execute_process("usc5-552-doj")
    execute_process("nodejs-governance")
    workbook_roundtrip("usc5-552-doj")
    attendance_chart()
    coverage_matrix()
    never_discussed("nodejs-governance")
    soundness_discovered("receipt-xes")
    soundness_from_order("usc5-552-doj")
    soundness_from_order("nodejs-governance")
    variants("receipt-xes")
    rdf("usc5-552-doj")
    rdf("nodejs-tsc-charter")
    graph_db("usc5-552-doj")
    site("usc5-552-doj")
    document_roundtrip("usc5-552-doj")
    bpmn_3d("usc5-552-doj")
    bpmn_3d("nodejs-governance")
    policy_page("majority")
    policy_page("usc5-552-doj")
    heatmap("receipt-csv")
    who_must_what("usc5-552-doj")
    timeline("nodejs-governance")
    bottlenecks("receipt-xes")
    at_risk("receipt-csv")
    checklist("nodejs-tsc-charter")
    checklist("usc5-552-doj")
    cross_document(["nodejs-tsc-charter", "nodejs-governance"])
    site_roundtrip("usc5-552-doj")
    seal_all()
    ITIL_REPORT.update(run_itil())
    write_md()
    print(f"proofs {len(PROOFS)}; itil deliverables {sum(len(p['deliverables']) for p in ITIL_REPORT.get('practices', []))}; wrote {MD.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
